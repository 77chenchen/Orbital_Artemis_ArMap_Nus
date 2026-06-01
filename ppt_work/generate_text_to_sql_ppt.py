from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


OUT = Path(__file__).resolve().parent / "基于LLM的Text-to-SQL技术演进_从模式映射到SchemaLinking.pptx"

W, H = Cm(33.867), Cm(19.05)

NAVY = RGBColor(18, 32, 47)
INK = RGBColor(36, 43, 52)
MUTED = RGBColor(91, 104, 118)
SUBTLE = RGBColor(130, 143, 156)
BLUE = RGBColor(36, 107, 201)
CYAN = RGBColor(34, 159, 174)
GREEN = RGBColor(50, 143, 96)
AMBER = RGBColor(218, 151, 45)
RED = RGBColor(192, 76, 76)
PURPLE = RGBColor(116, 91, 189)
PALE = RGBColor(246, 249, 252)
PALE_BLUE = RGBColor(235, 244, 255)
PALE_GREEN = RGBColor(237, 248, 242)
PALE_AMBER = RGBColor(255, 247, 231)
PALE_RED = RGBColor(253, 239, 239)
LINE = RGBColor(211, 221, 230)
WHITE = RGBColor(255, 255, 255)

FONT = "PingFang SC"
MONO = "Menlo"

SOURCE_SHORT = "资料来源：Spider、WikiSQL、RAT-SQL、PICARD、BIRD、Spider 2.0、DIN-SQL、DAIL-SQL、CHESS 等论文与 benchmark；整理见末页。"


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=None, font=FONT, line_spacing=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.line_spacing = line_spacing
    if align:
        p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_bg(slide, color=RGBColor(250, 252, 254)):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_footer(slide, page_no=None, source=SOURCE_SHORT):
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.05), Cm(17.86), Cm(31.8), Cm(0.02))
    set_fill(rule, LINE)
    rule.line.fill.background()
    add_text(slide, Cm(1.05), Cm(18.03), Cm(29.0), Cm(0.62), source, size=6.4, color=RGBColor(110, 122, 134))
    if page_no is not None:
        add_text(slide, Cm(31.55), Cm(18.0), Cm(1.2), Cm(0.52), f"{page_no:02d}", size=8.5, color=SUBTLE, align=PP_ALIGN.RIGHT)


def add_title(slide, title, kicker=None):
    if kicker:
        add_text(slide, Cm(1.25), Cm(0.72), Cm(13.5), Cm(0.52), kicker, size=10.5, color=BLUE, bold=True)
    add_text(slide, Cm(1.18), Cm(1.18), Cm(27.5), Cm(1.05), title, size=27, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.22), Cm(2.42), Cm(2.25), Cm(0.08))
    set_fill(line, CYAN)
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent=BLUE, fill=WHITE, body_size=11.6, title_size=15.2):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    set_fill(card, fill)
    set_line(card, RGBColor(220, 228, 236))
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Cm(0.12), h)
    set_fill(bar, accent)
    bar.line.fill.background()
    add_text(slide, x + Cm(0.42), y + Cm(0.34), w - Cm(0.72), Cm(0.64), title, size=title_size, color=NAVY, bold=True)
    add_text(slide, x + Cm(0.42), y + Cm(1.12), w - Cm(0.72), h - Cm(1.28), body, size=body_size, color=MUTED, line_spacing=1.08)
    return card


def add_pill(slide, x, y, w, text, color=BLUE, text_color=WHITE):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Cm(0.65))
    pill.adjustments[0] = 0.35
    set_fill(pill, color)
    pill.line.fill.background()
    add_text(slide, x + Cm(0.10), y + Cm(0.13), w - Cm(0.20), Cm(0.35), text, size=9.5, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return pill


def add_bullets(slide, x, y, w, h, bullets, size=13.4, color=INK, gap=3):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(136, 154, 171), width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_line(slide, x1, y1, x2, y2, color=RGBColor(136, 154, 171), width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_metric(slide, x, y, w, h, value, label, color=BLUE, note=None):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    set_fill(box, WHITE)
    set_line(box)
    add_text(slide, x + Cm(0.25), y + Cm(0.28), w - Cm(0.5), Cm(0.86), value, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Cm(0.3), y + Cm(1.22), w - Cm(0.6), Cm(0.58), label, size=9.8, color=MUTED, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, x + Cm(0.3), y + Cm(1.84), w - Cm(0.6), Cm(0.52), note, size=7.5, color=SUBTLE, align=PP_ALIGN.CENTER)
    return box


def add_stage(slide, x, y, w, h, title, body, color):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    set_fill(box, WHITE)
    set_line(box, RGBColor(215, 225, 234))
    cap = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Cm(0.82))
    set_fill(cap, color)
    cap.line.fill.background()
    add_text(slide, x + Cm(0.18), y + Cm(0.19), w - Cm(0.36), Cm(0.38), title, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Cm(0.34), y + Cm(1.12), w - Cm(0.68), h - Cm(1.28), body, size=10.8, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.04)
    return box


def add_table(slide, x, y, col_widths, row_h, headers, rows, header_color=NAVY, font_size=9.5):
    total_w = sum(col_widths)
    h = row_h * (len(rows) + 1)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, total_w, h)
    bg.adjustments[0] = 0.04
    set_fill(bg, WHITE)
    set_line(bg, RGBColor(219, 227, 235))
    cx = x
    for idx, head in enumerate(headers):
        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cx, y, col_widths[idx], row_h)
        set_fill(cell, header_color)
        cell.line.color.rgb = header_color
        add_text(slide, cx + Cm(0.10), y + Cm(0.17), col_widths[idx] - Cm(0.2), row_h - Cm(0.18), head, size=9.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        cx += col_widths[idx]
    for r, row in enumerate(rows):
        cy = y + row_h * (r + 1)
        cx = x
        for c, txt in enumerate(row):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cx, cy, col_widths[c], row_h)
            set_fill(cell, RGBColor(252, 254, 255) if r % 2 == 0 else RGBColor(247, 250, 253))
            set_line(cell, RGBColor(226, 233, 240), width=0.6)
            add_text(slide, cx + Cm(0.14), cy + Cm(0.12), col_widths[c] - Cm(0.28), row_h - Cm(0.18), txt, size=font_size, color=INK if c == 0 else MUTED, bold=(c == 0), line_spacing=1.0)
            cx += col_widths[c]


def add_code_box(slide, x, y, w, h, code, accent=CYAN):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.05
    set_fill(box, RGBColor(26, 38, 50))
    box.line.fill.background()
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Cm(0.62))
    set_fill(top, RGBColor(38, 54, 69))
    top.line.fill.background()
    for i, color in enumerate([RED, AMBER, GREEN]):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Cm(0.34 + i * 0.36), y + Cm(0.22), Cm(0.13), Cm(0.13))
        set_fill(dot, color)
        dot.line.fill.background()
    add_text(slide, x + Cm(0.36), y + Cm(0.90), w - Cm(0.72), h - Cm(1.1), code, size=10.3, color=RGBColor(226, 238, 245), font=MONO, line_spacing=1.08)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y + h - Cm(0.08), w, Cm(0.08))
    set_fill(strip, accent)
    strip.line.fill.background()


def cover(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, RGBColor(247, 250, 253))
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(10.35), H)
    set_fill(left, NAVY)
    left.line.fill.background()
    for i in range(9):
        y = Cm(3.0 + i * 1.35)
        add_line(s, Cm(1.25), y, Cm(8.75), y, RGBColor(52, 76, 99), width=0.7)
    for i, (tx, ty) in enumerate([(2.0, 3.15), (5.8, 4.5), (3.6, 5.85), (7.0, 7.2), (2.6, 8.55), (6.4, 9.9), (4.2, 11.25)]):
        n = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(tx), Cm(ty), Cm(0.42), Cm(0.42))
        set_fill(n, [CYAN, BLUE, GREEN, AMBER, PURPLE, RED, CYAN][i])
        n.line.fill.background()
    add_text(s, Cm(1.2), Cm(1.12), Cm(5.8), Cm(0.55), "技术综述 / 研究汇报", size=12, color=RGBColor(203, 222, 239), bold=True)
    add_text(s, Cm(12.0), Cm(4.1), Cm(18.8), Cm(2.4), "基于 LLM 的\nText-to-SQL 技术演进", size=36, color=NAVY, bold=True, line_spacing=0.95)
    add_text(s, Cm(12.05), Cm(7.25), Cm(17.5), Cm(0.92), "从模式映射到 Schema Linking 的深度解析", size=18, color=BLUE, bold=True)
    add_text(s, Cm(12.05), Cm(9.0), Cm(17.7), Cm(1.2), "核心判断：Text-to-SQL 的瓶颈已经从“会不会写 SQL”转向“能否把自然语言意图稳定绑定到正确的数据语义”。", size=15, color=MUTED)
    add_metric(s, Cm(12.1), Cm(11.35), Cm(4.6), Cm(2.35), "26", "页内容", BLUE)
    add_metric(s, Cm(17.05), Cm(11.35), Cm(4.6), Cm(2.35), "4", "技术阶段", GREEN)
    add_metric(s, Cm(22.0), Cm(11.35), Cm(4.6), Cm(2.35), "10+", "代表论文", AMBER)
    add_metric(s, Cm(26.95), Cm(11.35), Cm(4.6), Cm(2.35), "2026", "资料整理", PURPLE)
    add_text(s, Cm(12.1), Cm(15.35), Cm(15), Cm(0.6), "适合课程展示、组会汇报与技术方案讨论", size=12.5, color=SUBTLE)
    add_footer(s, page, "资料来源：公开论文、benchmark 页面与作者项目页；详细清单见末页。")


def thesis(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "一页结论：Text-to-SQL 的主线不是“生成”，而是“对齐”", "核心观点")
    add_text(s, Cm(1.35), Cm(3.05), Cm(29.6), Cm(0.85), "把自然语言变成 SQL 只是表象；真正困难的是在开放、异构、含业务语义的数据库中，把“词、值、列、表、关系、约束”绑定为可执行、可解释的查询计划。", size=16, color=MUTED)
    cards = [
        ("阶段 1：模式映射", "通过字符串匹配、同义词词典、实体识别，把问题中的词映射到列名、表名或枚举值。可解释，但脆弱。", BLUE),
        ("阶段 2：结构化解码", "Seq2Seq 之外，开始利用 SQL sketch、语法约束和列注意力，降低 SQL 表达的排列歧义。", CYAN),
        ("阶段 3：显式 Schema Linking", "将数据库 schema 视为图，建模表-列、外键、问题 token 与 schema item 的关系。RAT-SQL 是关键里程碑。", GREEN),
        ("阶段 4：LLM 系统工程", "LLM 负责语义推理与候选生成，系统负责检索、裁剪、约束、执行验证和人类反馈。", AMBER),
    ]
    for i, (title, body, color) in enumerate(cards):
        add_card(s, Cm(1.35 + i * 7.75), Cm(4.65), Cm(7.15), Cm(4.6), title, body, color, body_size=10.8, title_size=13.6)
    add_text(s, Cm(1.35), Cm(10.55), Cm(8.6), Cm(0.7), "我的判断", size=18, color=NAVY, bold=True)
    add_card(s, Cm(1.35), Cm(11.55), Cm(30.8), Cm(4.3), "Schema Linking 是 Text-to-SQL 的“语义路由层”", "在小 schema 上，强模型可以直接把 SQL 写出来；在企业级数据库中，LLM 的上下文、幻觉、业务口径冲突会放大。可靠系统必须先缩小候选 schema、补足值证据、显式记录链接证据，再让模型生成 SQL。换言之，Schema Linking 不再只是模型内部特征，而是可观测、可调试、可治理的工程组件。", GREEN, PALE_GREEN, body_size=14.2, title_size=16)
    add_footer(s, page)


def agenda(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "汇报结构", "从概念到落地")
    items = [
        ("01", "问题定义", "Text-to-SQL 为什么难：自然语言、数据库结构、业务语义三方不一致"),
        ("02", "技术演进", "从模板/语义解析、WikiSQL、Spider，到 RAT-SQL 和预训练模型"),
        ("03", "Schema Linking 深挖", "链接对象、链接粒度、图关系、值证据与错误类型"),
        ("04", "LLM 时代", "Prompt、RAG、schema pruning、多候选、执行反馈和 agent workflow"),
        ("05", "评测与落地", "Spider/BIRD/Spider 2.0 的差异，以及生产系统设计原则"),
    ]
    for i, (num, title, body) in enumerate(items):
        y = Cm(3.25 + i * 2.55)
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(1.35), y, Cm(1.32), Cm(1.32))
        set_fill(circ, [BLUE, CYAN, GREEN, AMBER, PURPLE][i])
        circ.line.fill.background()
        add_text(s, Cm(1.48), y + Cm(0.32), Cm(1.05), Cm(0.4), num, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Cm(3.05), y + Cm(0.03), Cm(8.0), Cm(0.52), title, size=17, color=NAVY, bold=True)
        add_text(s, Cm(3.05), y + Cm(0.76), Cm(24.8), Cm(0.78), body, size=13.4, color=MUTED)
        if i < len(items) - 1:
            add_line(s, Cm(2.01), y + Cm(1.34), Cm(2.01), y + Cm(2.28), RGBColor(183, 198, 211), width=1.2)
    add_card(s, Cm(23.7), Cm(3.25), Cm(8.2), Cm(10.6), "贯穿问题", "同一句话中的“最高分课程”“活跃用户”“去年同期”在不同数据库里可能对应完全不同的字段、过滤条件、聚合口径和时间维度。Text-to-SQL 的最终目标不是生成看起来正确的 SQL，而是生成与业务语义一致、可验证、可复用的查询。", BLUE, PALE_BLUE, body_size=13.2)
    add_footer(s, page)


def evolution_timeline(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "技术演进时间线：每一次进步都在减少“语义漂移”", "Historical View")
    milestones = [
        ("规则/模板", "固定领域\n词典映射", "可控但难扩展", BLUE),
        ("WikiSQL", "单表大规模\n结构化解码", "降低 SQL 语法噪声", CYAN),
        ("Spider", "跨域多表\n复杂 SQL", "Schema 泛化成为核心", GREEN),
        ("RAT-SQL", "关系感知图\n显式 linking", "链接被结构化建模", AMBER),
        ("PLM + IR", "T5/BERT\nsketch/约束", "语言泛化与结构结合", PURPLE),
        ("LLM 系统", "RAG/Agent\n执行验证", "从模型题变成系统题", RED),
    ]
    y = Cm(8.0)
    add_line(s, Cm(2.0), y, Cm(31.0), y, RGBColor(180, 195, 210), width=2.0)
    for i, (title, top, bottom, color) in enumerate(milestones):
        x = Cm(2.0 + i * 5.75)
        dot = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - Cm(0.32), y - Cm(0.32), Cm(0.64), Cm(0.64))
        set_fill(dot, color)
        dot.line.fill.background()
        add_text(s, x - Cm(2.25), Cm(4.05), Cm(4.5), Cm(0.62), title, size=15.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x - Cm(2.25), Cm(5.0), Cm(4.5), Cm(1.25), top, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
        add_line(s, x, Cm(6.55), x, y - Cm(0.38), RGBColor(190, 204, 216), width=1.0)
        add_line(s, x, y + Cm(0.38), x, Cm(10.15), RGBColor(190, 204, 216), width=1.0)
        add_text(s, x - Cm(2.35), Cm(10.45), Cm(4.7), Cm(1.15), bottom, size=11.1, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_card(s, Cm(1.45), Cm(13.0), Cm(30.5), Cm(2.6), "贯穿主线", "从模板到 LLM，方法越来越开放，但系统可靠性仍取决于一个朴素问题：用户话里的概念到底应当落到哪些表、列、值、JOIN path 和业务规则上。", GREEN, PALE_GREEN, body_size=13.2)
    add_footer(s, page)


def task_definition(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Text-to-SQL 任务：从一句话到可执行查询", "问题定义")
    add_stage(s, Cm(1.35), Cm(4.25), Cm(5.6), Cm(3.0), "自然语言问题", "“列出 2024 年销售额最高的三个品类及同比增长率”", BLUE)
    add_stage(s, Cm(8.05), Cm(4.25), Cm(5.6), Cm(3.0), "Schema / Metadata", "tables, columns, types, PK/FK, comments, business glossary", CYAN)
    add_stage(s, Cm(14.75), Cm(4.25), Cm(5.6), Cm(3.0), "Schema Linking", "问题词、值、列、表、关系之间的显式对齐证据", GREEN)
    add_stage(s, Cm(21.45), Cm(4.25), Cm(5.6), Cm(3.0), "SQL 生成", "SELECT / JOIN / WHERE / GROUP BY / ORDER BY / dialect", AMBER)
    add_stage(s, Cm(28.15), Cm(4.25), Cm(4.2), Cm(3.0), "执行验证", "结果、代价、权限、安全", RED)
    for x1 in [Cm(6.95), Cm(13.65), Cm(20.35), Cm(27.05)]:
        add_arrow(s, x1, Cm(5.75), x1 + Cm(1.0), Cm(5.75))
    add_code_box(s, Cm(1.45), Cm(9.25), Cm(14.55), Cm(5.25), "SELECT c.category_name,\n       SUM(o.amount) AS sales,\n       (SUM(o.amount) - SUM(prev.amount)) / SUM(prev.amount) AS yoy\nFROM orders o\nJOIN categories c ON o.category_id = c.id\nLEFT JOIN orders prev ON ...\nWHERE o.order_date BETWEEN '2024-01-01' AND '2024-12-31'\nGROUP BY c.category_name\nORDER BY sales DESC\nLIMIT 3;")
    add_card(s, Cm(17.0), Cm(9.25), Cm(14.95), Cm(5.25), "关键不是 SQL 模板，而是语义绑定", "SQL 的语法空间有限，但数据库语义空间很大：字段名可能缩写，业务概念可能跨表，值可能需要模糊匹配，日期和指标口径需要上下文。Schema Linking 负责把开放语言压缩到可执行的数据库操作空间。", GREEN, PALE_GREEN, body_size=13.2)
    add_footer(s, page)


def hard_cases(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "为什么 Text-to-SQL 比“翻译”更难", "挑战拆解")
    data = [
        ("词面不一致", "“老师/讲师/Instructor”“销售额/GMV/revenue”经常不是字段名原文。", BLUE),
        ("值证据缺失", "用户问“新加坡校区”，表里可能是 campus_id=2 或 campus_name='Kent Ridge'。", CYAN),
        ("结构推理", "JOIN path、外键方向、桥表、多对多关系需要结构搜索。", GREEN),
        ("业务口径", "“活跃用户”“订单成功率”往往是组织内约定，不一定写在 schema。", AMBER),
        ("SQL 等价性", "不同 SQL 可产生同一结果，exact match 与真实正确性存在偏差。", PURPLE),
        ("安全与权限", "生产环境还要考虑注入、越权、长查询、敏感字段和审计。", RED),
    ]
    for i, (title, body, color) in enumerate(data):
        row, col = divmod(i, 3)
        add_card(s, Cm(1.35 + col * 10.35), Cm(3.55 + row * 5.6), Cm(9.55), Cm(4.65), title, body, color, body_size=12.0)
    add_text(s, Cm(1.4), Cm(15.55), Cm(29.5), Cm(0.85), "独立思考：模型越强，越容易掩盖链接错误。它会生成语法漂亮、执行成功、但业务语义错位的 SQL；因此链接证据和执行审计比单纯提升生成模型更重要。", size=14.6, color=NAVY, bold=True)
    add_footer(s, page)


def concept_mapping_linking(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "从 Schema Mapping 到 Schema Linking：概念边界", "核心概念")
    add_table(
        s,
        Cm(1.35),
        Cm(3.25),
        [Cm(5.2), Cm(8.0), Cm(8.0), Cm(8.0)],
        Cm(1.55),
        ["维度", "Schema Mapping", "Schema Linking", "Schema Grounding"],
        [
            ("目标", "词面或实体到 schema item 的候选映射", "问题 token / 值 / 表列 / 关系的结构化对齐", "把链接证据落到真实数据、权限、业务定义和执行环境"),
            ("典型方法", "字符串匹配、同义词、NER、列名相似度", "关系图、注意力、显式链接标签、schema pruning", "RAG、元数据检索、值采样、执行反馈、数据目录"),
            ("优点", "简单、可解释、低成本", "能建模跨表关系和组合语义", "更贴近生产，能处理脏数据和业务口径"),
            ("短板", "词面变化大时脆弱", "链接错误会传导到生成器", "工程复杂、依赖治理质量与权限体系"),
        ],
        font_size=9.1,
    )
    add_card(s, Cm(1.55), Cm(11.8), Cm(30.4), Cm(3.9), "一个清晰区分", "Mapping 解决“可能是哪一列”；Linking 解决“这些列和值如何组成一个查询意图”；Grounding 解决“这个意图在当前组织的数据资产和业务规则中是否成立”。LLM 时代真正缺的通常不是语言理解，而是 Grounding 能力。", GREEN, PALE_GREEN, body_size=14.0)
    add_footer(s, page)


def pre_neural(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "早期路线：规则、模板与语义解析", "阶段 0")
    add_text(s, Cm(1.35), Cm(3.1), Cm(29.2), Cm(0.8), "早期系统强调可控性：限制问题表达、维护词典和模板，用规则把自然语言解析为逻辑形式或 SQL 片段。", size=15.5, color=MUTED)
    add_stage(s, Cm(2.0), Cm(5.0), Cm(6.0), Cm(3.2), "Lexicon", "业务词典\n同义词\n实体列表", BLUE)
    add_stage(s, Cm(9.4), Cm(5.0), Cm(6.0), Cm(3.2), "Template", "固定句式\nSQL slots\n规则组合", CYAN)
    add_stage(s, Cm(16.8), Cm(5.0), Cm(6.0), Cm(3.2), "Semantic Parser", "lambda calculus\nSQL AST\n约束搜索", GREEN)
    add_stage(s, Cm(24.2), Cm(5.0), Cm(6.0), Cm(3.2), "Database", "执行\n返回结果\n人工修正", AMBER)
    for x in [Cm(8.0), Cm(15.4), Cm(22.8)]:
        add_arrow(s, x, Cm(6.6), x + Cm(1.35), Cm(6.6))
    add_card(s, Cm(1.35), Cm(10.1), Cm(9.6), Cm(4.7), "优势", "高可解释、低幻觉、容易嵌入权限与审计；在窄领域、强约束场景仍然有效。", GREEN, PALE_GREEN)
    add_card(s, Cm(12.1), Cm(10.1), Cm(9.6), Cm(4.7), "瓶颈", "规则维护成本高；跨域泛化差；遇到长问题、隐含条件和复杂 JOIN 时迅速变脆。", RED, PALE_RED)
    add_card(s, Cm(22.85), Cm(10.1), Cm(9.1), Cm(4.7), "遗产", "今天的 LLM 系统仍在复用这条思路：词典、元数据、业务 glossary、SQL 模板都变成检索上下文。", BLUE, PALE_BLUE)
    add_footer(s, page)


def neural_wikisql(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "WikiSQL 时代：从 Seq2Seq 到结构化解码", "2017-2018")
    add_text(s, Cm(1.35), Cm(3.1), Cm(29.5), Cm(0.8), "WikiSQL 把任务推向大规模训练；Seq2SQL、SQLNet 等方法发现：SQL 不是普通句子，直接序列化会带来排列歧义和无效查询。", size=15.2, color=MUTED)
    add_card(s, Cm(1.35), Cm(4.65), Cm(9.6), Cm(5.2), "Seq2SQL", "引入强化学习优化执行结果，处理 WHERE 条件顺序等 SQL 等价问题；但训练不稳定，结构约束有限。", BLUE, PALE_BLUE, body_size=12.0)
    add_card(s, Cm(12.1), Cm(4.65), Cm(9.6), Cm(5.2), "SQLNet", "采用 sketch 与 sequence-to-set，避免 order-matters；列注意力让模型聚焦 schema item。", GREEN, PALE_GREEN, body_size=12.0)
    add_card(s, Cm(22.85), Cm(4.65), Cm(9.1), Cm(5.2), "TypeSQL", "用列类型、值类型和知识增强，开始把“问题词是否像一个值/列”显式编码。", AMBER, PALE_AMBER, body_size=12.0)
    add_text(s, Cm(1.5), Cm(11.55), Cm(7.4), Cm(0.55), "这一阶段的隐含假设", size=16, color=NAVY, bold=True)
    add_bullets(s, Cm(1.55), Cm(12.5), Cm(29.8), Cm(3.3), [
        "多为单表或结构相对简单，核心难点是 SELECT/WHERE/聚合槽位预测。",
        "Schema Linking 主要体现为列注意力和类型匹配，还没有成为独立模块。",
        "对企业数据库中跨表 JOIN、业务指标、脏值、注释元数据的挑战覆盖不足。",
    ], size=14.4)
    add_footer(s, page, "资料来源：Seq2SQL、SQLNet、TypeSQL、WikiSQL 相关论文；详细清单见末页。")


def spider_shift(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Spider：把问题推向“跨域、多表、复杂 SQL”", "2018 之后的转折")
    add_metric(s, Cm(1.45), Cm(3.55), Cm(5.6), Cm(2.45), "10,181", "questions", BLUE)
    add_metric(s, Cm(7.65), Cm(3.55), Cm(5.6), Cm(2.45), "5,693", "unique SQL", CYAN)
    add_metric(s, Cm(13.85), Cm(3.55), Cm(5.6), Cm(2.45), "200", "databases", GREEN)
    add_metric(s, Cm(20.05), Cm(3.55), Cm(5.6), Cm(2.45), "138", "domains", AMBER)
    add_metric(s, Cm(26.25), Cm(3.55), Cm(5.6), Cm(2.45), "12.4%", "早期最佳 EM", RED, note="database split")
    add_card(s, Cm(1.35), Cm(7.2), Cm(14.6), Cm(5.05), "范式变化", "训练和测试数据库不同，模型不能靠记忆字段或 SQL 模板；必须泛化到未见 schema，并理解表间关系。", BLUE, PALE_BLUE, body_size=13.2)
    add_card(s, Cm(17.05), Cm(7.2), Cm(14.9), Cm(5.05), "对 Schema Linking 的影响", "schema item 的候选空间变大，列名不再直接出现在问题里，JOIN path 与外键结构成为关键。Schema Linking 从辅助特征变成主战场。", GREEN, PALE_GREEN, body_size=13.2)
    add_text(s, Cm(1.5), Cm(13.65), Cm(29.7), Cm(1.0), "Spider 的价值不只是一个数据集，而是把 Text-to-SQL 从“在固定数据库上拟合语言模式”改造成“在新数据库上进行结构泛化”。", size=16.2, color=NAVY, bold=True)
    add_footer(s, page, "资料来源：Spider paper, arXiv:1809.08887。")


def taxonomy(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Schema Linking 到底链接什么", "链接对象")
    items = [
        ("Token ↔ Column", "“价格” ↔ price / amount / revenue", BLUE),
        ("Token ↔ Table", "“学生” ↔ student / enrollment", CYAN),
        ("Value ↔ Cell", "“CS101” ↔ course_code='CS101'", GREEN),
        ("Column ↔ Column", "FK/PK、join path、桥表", AMBER),
        ("Phrase ↔ Operator", "“最多/平均/去年” ↔ ORDER BY / AVG / time filter", PURPLE),
        ("Intent ↔ Business Rule", "“活跃” ↔ 最近 30 天登录且未注销", RED),
    ]
    center = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(13.2), Cm(6.0), Cm(7.6), Cm(3.6))
    set_fill(center, NAVY)
    center.line.fill.background()
    add_text(s, Cm(14.0), Cm(6.9), Cm(6.0), Cm(0.8), "Schema Linking", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    positions = [(1.4, 3.25), (12.1, 3.25), (22.8, 3.25), (1.4, 11.3), (12.1, 11.3), (22.8, 11.3)]
    for (title, body, color), (x, y) in zip(items, positions):
        add_card(s, Cm(x), Cm(y), Cm(9.15), Cm(3.6), title, body, color, body_size=12.3)
        add_arrow(s, Cm(x + 4.58), Cm(y + (3.6 if y < 6 else 0)), Cm(17.0), Cm(7.8), RGBColor(170, 184, 197), width=1.0)
    add_footer(s, page)


def rat_sql(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "RAT-SQL：把 Linking 做成关系感知图编码", "代表方法")
    add_text(s, Cm(1.35), Cm(3.0), Cm(29.5), Cm(0.78), "RAT-SQL 将问题 token、表、列共同放入图结构，用 relation-aware self-attention 编码关系类型，同时融合 schema encoding 与 schema linking。", size=15.2, color=MUTED)
    nodes = [
        ("Q1", "show", 3.2, 6.2, BLUE), ("Q2", "students", 5.0, 4.8, BLUE), ("Q3", "courses", 7.0, 6.3, BLUE),
        ("T", "student", 13.4, 4.9, CYAN), ("C", "student.name", 16.0, 3.8, GREEN), ("C", "student.id", 16.2, 6.0, GREEN),
        ("T", "enrollment", 22.0, 5.1, CYAN), ("C", "course_id", 24.8, 4.0, GREEN), ("C", "student_id", 24.8, 6.2, GREEN),
    ]
    for _, _, x1, y1, _ in nodes:
        for _, _, x2, y2, _ in nodes:
            if abs(x1 - x2) > 6 and abs(y1 - y2) < 2 and x1 < x2:
                add_line(s, Cm(x1 + 0.5), Cm(y1 + 0.35), Cm(x2), Cm(y2 + 0.35), RGBColor(218, 226, 234), width=0.8)
    for kind, label, x, y, color in nodes:
        shp = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(2.5), Cm(0.78))
        shp.adjustments[0] = 0.18
        set_fill(shp, color)
        shp.line.fill.background()
        add_text(s, Cm(x + 0.1), Cm(y + 0.17), Cm(2.3), Cm(0.35), label, size=8.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_card(s, Cm(1.35), Cm(9.6), Cm(9.7), Cm(4.8), "关系类型", "table-column、primary key、foreign key、same table、question-schema exact/partial/value match 等。", BLUE, PALE_BLUE, body_size=12.4)
    add_card(s, Cm(12.1), Cm(9.6), Cm(9.7), Cm(4.8), "技术贡献", "不是把 schema 当普通文本拼接，而是显式建模数据库结构与语言对齐。", GREEN, PALE_GREEN, body_size=12.4)
    add_card(s, Cm(22.85), Cm(9.6), Cm(9.1), Cm(4.8), "局限", "依赖监督数据与静态 schema；面对超大 schema、业务 glossary 和真实值脏数据仍不够。", RED, PALE_RED, body_size=12.4)
    add_footer(s, page, "资料来源：RAT-SQL paper, arXiv:1911.04942。")


def intermediate_ir(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "中间表示：让模型先想结构，再填 Schema", "IR / Sketch 路线")
    add_text(s, Cm(1.35), Cm(3.0), Cm(29.5), Cm(0.78), "IRNet、SemQL、NatSQL 等思路把 SQL 的生成拆成“抽象结构 + schema item 填充”，降低 SQL 语法和数据库绑定同时学习的难度。", size=15.2, color=MUTED)
    add_stage(s, Cm(2.0), Cm(5.0), Cm(6.2), Cm(3.4), "Question", "自然语言\n含省略与歧义", BLUE)
    add_stage(s, Cm(9.4), Cm(5.0), Cm(6.2), Cm(3.4), "IR / Sketch", "SELECT-AGG\nFILTER\nJOIN slots", CYAN)
    add_stage(s, Cm(16.8), Cm(5.0), Cm(6.2), Cm(3.4), "Schema Filling", "列/表/值链接\n约束检查", GREEN)
    add_stage(s, Cm(24.2), Cm(5.0), Cm(6.2), Cm(3.4), "SQL", "方言化\n执行验证", AMBER)
    for x in [Cm(8.2), Cm(15.6), Cm(23.0)]:
        add_arrow(s, x, Cm(6.75), x + Cm(1.2), Cm(6.75))
    add_card(s, Cm(1.45), Cm(10.25), Cm(14.55), Cm(4.5), "为什么有效", "复杂 SQL 的结构语义与 schema 名称绑定是两类问题。先预测结构，可以减少解码空间；再做 schema filling，可以更清楚地定位链接错误。", GREEN, PALE_GREEN, body_size=13.2)
    add_card(s, Cm(17.0), Cm(10.25), Cm(14.95), Cm(4.5), "对 LLM 的启发", "今天的 prompt decomposition、planner-generator、agent routing 本质上延续了 IR 路线：先拆解任务和查询骨架，再选择 schema 与生成 SQL。", BLUE, PALE_BLUE, body_size=13.2)
    add_footer(s, page, "资料来源：IRNet / SemQL、NatSQL、SQL sketch 相关论文。")


def plm_resdsql(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "预训练模型时代：把 Schema 当文本还是当结构", "2019-2023")
    add_table(
        s,
        Cm(1.35),
        Cm(3.35),
        [Cm(5.1), Cm(7.9), Cm(8.3), Cm(8.2)],
        Cm(1.55),
        ["方法方向", "做法", "Schema Linking 角色", "典型收益/风险"],
        [
            ("BRIDGE", "把问题、表、列和值序列化输入 PLM", "用文本上下文与值证据增强链接", "实现简单，但长 schema 易超上下文"),
            ("Graphix-T5", "T5 + 图结构层", "保留 schema 图关系", "兼顾预训练与结构，但系统复杂"),
            ("RESDSQL", "解耦 schema linking 与 skeleton parsing", "先排序/裁剪相关 schema，再生成骨架", "降低干扰；错误裁剪会漏答案"),
            ("PICARD", "增量解析约束解码", "主要保证语法合法，间接降低无效 SQL", "不能保证语义链接正确"),
        ],
        font_size=8.8,
    )
    add_card(s, Cm(1.45), Cm(11.45), Cm(30.5), Cm(3.8), "阶段性结论", "预训练模型提升了语言泛化，但没有消除 Schema Linking。相反，随着 schema 被文本化塞进模型，如何选择、排序、压缩、标注 schema 变成决定 token 效率和准确率的关键。", BLUE, PALE_BLUE, body_size=14)
    add_footer(s, page, "资料来源：BRIDGE、Graphix-T5、RESDSQL、PICARD 等论文；详细清单见末页。")


def constrained_feedback(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "约束解码与执行反馈：生成之后还要“验算”", "可靠性工具箱")
    add_text(s, Cm(1.35), Cm(3.0), Cm(29.5), Cm(0.78), "Text-to-SQL 的正确性可以拆成三层：语法合法、可执行、语义正确。PICARD 和 execution-guided decoding 主要覆盖前两层，第三层仍需要链接证据和业务验证。", size=15.0, color=MUTED)
    add_stage(s, Cm(1.55), Cm(5.15), Cm(6.0), Cm(3.2), "1. 语法约束", "SQL grammar\nAST parser\ninadmissible token rejection", BLUE)
    add_stage(s, Cm(9.15), Cm(5.15), Cm(6.0), Cm(3.2), "2. 类型约束", "列类型\n聚合函数\n比较操作", CYAN)
    add_stage(s, Cm(16.75), Cm(5.15), Cm(6.0), Cm(3.2), "3. 执行反馈", "error message\nempty result\nquery cost", GREEN)
    add_stage(s, Cm(24.35), Cm(5.15), Cm(6.0), Cm(3.2), "4. 语义验收", "unit tests\nbusiness rule\nhuman review", AMBER)
    for x in [Cm(7.55), Cm(15.15), Cm(22.75)]:
        add_arrow(s, x, Cm(6.75), x + Cm(1.55), Cm(6.75))
    add_card(s, Cm(1.55), Cm(10.45), Cm(14.5), Cm(4.25), "PICARD 的启示", "语言模型可以被语法解析器“拦住手”：每一步只允许产生仍可被解析为 SQL 的 token，从而显著减少无效输出。", GREEN, PALE_GREEN, body_size=13.1)
    add_card(s, Cm(17.05), Cm(10.45), Cm(14.9), Cm(4.25), "但还不够", "能执行不代表问对了表、连对了字段。语义正确性需要问题意图、schema 链接证据、样例结果和业务口径共同校验。", RED, PALE_RED, body_size=13.1)
    add_footer(s, page, "资料来源：PICARD paper, arXiv:2109.05093；execution-guided Text-to-SQL 相关研究。")


def llm_paradigm(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "LLM 时代：从“训练一个解析器”到“搭一个查询系统”", "2023 之后")
    add_text(s, Cm(1.35), Cm(3.0), Cm(29.5), Cm(0.78), "ChatGPT/Codex 一类模型证明了通用语言模型具备强 SQL 先验，但真实表现高度依赖上下文组织、示例选择、schema 选择与验证流程。", size=15.0, color=MUTED)
    steps = [
        ("Metadata RAG", "表说明、列注释、业务 glossary、历史查询", BLUE),
        ("Schema Pruning", "从数千列缩小到几十列候选", CYAN),
        ("Prompt Planning", "问题分解、SQL sketch、few-shot 示例", GREEN),
        ("Candidate Generation", "多路径生成、self-consistency、温度采样", AMBER),
        ("Verifier", "执行、单元测试、代价、权限、安全", RED),
    ]
    x = Cm(1.35)
    for i, (title, body, color) in enumerate(steps):
        add_stage(s, x + Cm(i * 6.15), Cm(5.0), Cm(5.45), Cm(3.6), title, body, color)
        if i < len(steps) - 1:
            add_arrow(s, x + Cm(i * 6.15 + 5.45), Cm(6.8), x + Cm(i * 6.15 + 6.15), Cm(6.8))
    add_card(s, Cm(1.45), Cm(10.15), Cm(9.7), Cm(4.95), "模型能力", "理解自然语言、解释 schema 注释、生成 SQL 草案、修复执行错误。", BLUE, PALE_BLUE, body_size=12.6)
    add_card(s, Cm(12.1), Cm(10.15), Cm(9.7), Cm(4.95), "系统能力", "检索、裁剪、约束、缓存、日志、审计、权限和回归测试。", GREEN, PALE_GREEN, body_size=12.6)
    add_card(s, Cm(22.75), Cm(10.15), Cm(9.2), Cm(4.95), "核心平衡", "更长上下文不等于更强链接；精准证据通常比把整个数据库塞进去更重要。", AMBER, PALE_AMBER, body_size=12.6)
    add_footer(s, page)


def prompt_methods(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Prompt Engineering：把 Schema Linking 外显到提示词", "LLM 方法一")
    add_table(
        s,
        Cm(1.35),
        Cm(3.2),
        [Cm(5.5), Cm(8.2), Cm(8.2), Cm(7.6)],
        Cm(1.48),
        ["策略", "做法", "价值", "风险"],
        [
            ("Schema Serialization", "DDL、列名、类型、外键、样例值、注释", "给模型数据库语义", "上下文膨胀、噪声干扰"),
            ("Few-shot Selection", "按问题/schema 相似度选示例", "迁移 SQL 模式和链接模式", "示例偏差会诱导错误"),
            ("Decomposition", "分类、schema linking、SQL 生成、自校正", "降低单步推理难度", "多步误差累积"),
            ("Self-correction", "执行报错后让模型修复", "提高可执行率", "容易修语法不修语义"),
        ],
        font_size=8.9,
    )
    add_card(s, Cm(1.4), Cm(10.75), Cm(14.7), Cm(4.7), "DIN-SQL", "把任务分为 schema linking、query classification、SQL generation、self-correction 等子任务，说明 LLM 需要显式推理脚手架。", GREEN, PALE_GREEN, body_size=13.0)
    add_card(s, Cm(17.0), Cm(10.75), Cm(14.95), Cm(4.7), "DAIL-SQL", "系统比较问题表示、示例选择、示例组织，强调 prompt token efficiency：不是越多越好，而是越相关越好。", BLUE, PALE_BLUE, body_size=13.0)
    add_footer(s, page, "资料来源：DIN-SQL arXiv:2304.11015；DAIL-SQL / benchmark evaluation arXiv:2308.15363。")


def rag_pruning(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "RAG 与 Schema Pruning：大库里的第一性问题", "LLM 方法二")
    add_text(s, Cm(1.35), Cm(3.0), Cm(29.5), Cm(0.78), "当数据库有数百张表、数千列时，真正的任务首先是找上下文。Schema Linking 从模型内部注意力变成检索与排序问题。", size=15.0, color=MUTED)
    add_stage(s, Cm(1.6), Cm(5.0), Cm(5.3), Cm(3.2), "Query", "用户问题\n会话上下文", BLUE)
    add_stage(s, Cm(8.25), Cm(5.0), Cm(5.3), Cm(3.2), "Retriever", "BM25 / embedding\n业务词典", CYAN)
    add_stage(s, Cm(14.9), Cm(5.0), Cm(5.3), Cm(3.2), "Ranker", "表列相关性\nJOIN 可达性", GREEN)
    add_stage(s, Cm(21.55), Cm(5.0), Cm(5.3), Cm(3.2), "Sub-schema", "候选表列\n样例值", AMBER)
    add_stage(s, Cm(28.2), Cm(5.0), Cm(3.65), Cm(3.2), "LLM", "生成\n验证", RED)
    for x in [Cm(6.9), Cm(13.55), Cm(20.2), Cm(26.85)]:
        add_arrow(s, x, Cm(6.6), x + Cm(1.25), Cm(6.6))
    add_card(s, Cm(1.45), Cm(10.05), Cm(9.7), Cm(4.85), "裁剪的收益", "降低 token 成本；减少无关字段干扰；让模型更专注于少量可验证证据。", GREEN, PALE_GREEN, body_size=12.8)
    add_card(s, Cm(12.1), Cm(10.05), Cm(9.7), Cm(4.85), "裁剪的风险", "一旦漏掉关键表/列，后续生成再强也难以恢复；需要召回优先、精排其次。", RED, PALE_RED, body_size=12.8)
    add_card(s, Cm(22.75), Cm(10.05), Cm(9.2), Cm(4.85), "工程建议", "保留 top-k 证据、链接分数、JOIN path；允许模型请求补充 schema，而不是一次性闭卷。", BLUE, PALE_BLUE, body_size=12.8)
    add_footer(s, page)


def multi_agent(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "多 Agent / 多候选：用测试时计算换可靠性", "LLM 方法三")
    add_text(s, Cm(1.35), Cm(3.0), Cm(29.5), Cm(0.78), "CHESS、MAC-SQL、CHASE-SQL 等方法把 Text-to-SQL 拆成多个角色：检索、schema 选择、候选生成、执行验证、偏好选择。", size=15.0, color=MUTED)
    roles = [
        ("信息检索器", "找相关表列、样例值、文档", BLUE),
        ("Schema Selector", "压缩 schema，保留 join path", CYAN),
        ("SQL Generator", "生成多条候选 SQL", GREEN),
        ("Unit Tester", "构造自然语言测试或执行检查", AMBER),
        ("Selector / Refiner", "比较结果、修复错误、选最终答案", RED),
    ]
    for i, (title, body, color) in enumerate(roles):
        add_stage(s, Cm(1.35 + i * 6.15), Cm(5.0), Cm(5.45), Cm(3.8), title, body, color)
        if i < len(roles) - 1:
            add_arrow(s, Cm(6.8 + i * 6.15), Cm(6.9), Cm(7.5 + i * 6.15), Cm(6.9))
    add_card(s, Cm(1.45), Cm(10.35), Cm(14.55), Cm(4.55), "何时值得多 Agent", "问题复杂、schema 很大、结果需要高置信度，或者有足够执行环境和预算时，多候选 + 验证比单次生成更稳。", GREEN, PALE_GREEN, body_size=13.1)
    add_card(s, Cm(17.0), Cm(10.35), Cm(14.95), Cm(4.55), "何时不值得", "低延迟、低成本、窄领域、高频简单查询，过多 agent 会增加成本、状态管理和调试复杂度。", AMBER, PALE_AMBER, body_size=13.1)
    add_footer(s, page, "资料来源：CHESS arXiv:2405.16755；MAC-SQL arXiv:2312.11242；CHASE-SQL arXiv:2410.01943。")


def benchmarks(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Benchmark 演进：评测压力越来越接近真实世界", "评测体系")
    rows = [
        ("WikiSQL", "单表、较简单 SQL", "结构化解码、列注意力", "覆盖复杂 JOIN 不足"),
        ("Spider", "跨域、多表、复杂 SQL", "schema 泛化、结构链接", "数据规模和真实值复杂度有限"),
        ("BIRD", "大数据库、真实值、37 专业领域", "值理解、外部知识、SQL 效率", "更贴近生产，但仍是任务集"),
        ("Spider 2.0", "企业级工作流、千列级 schema、多方言", "长上下文、文档检索、代码/环境交互", "对 agent 和工程系统提出新要求"),
    ]
    add_table(s, Cm(1.35), Cm(3.25), [Cm(4.5), Cm(8.2), Cm(8.6), Cm(8.2)], Cm(1.65), ["Benchmark", "特点", "推动的能力", "局限/启示"], rows, font_size=8.8)
    add_text(s, Cm(1.45), Cm(11.5), Cm(30), Cm(0.8), "注意：Leaderboards 不等于生产可用。企业落地还要评估权限、数据新鲜度、成本、延迟、可解释性、审计和错误恢复。", size=15.4, color=NAVY, bold=True)
    add_card(s, Cm(1.45), Cm(12.8), Cm(30.5), Cm(2.6), "Spider 2.0 的信号", "真实企业任务常需要查文档、理解项目代码、处理多 SQL 和方言差异；这意味着 Text-to-SQL 正在从单轮生成题变成数据工程 agent 题。", PURPLE, RGBColor(244, 241, 252), body_size=13.1)
    add_footer(s, page, "资料来源：WikiSQL、Spider arXiv:1809.08887；BIRD arXiv:2305.03111；Spider 2.0 arXiv:2411.07763。")


def errors(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "错误类型：Schema Linking 错了，SQL 往往看起来也很对", "失败模式")
    add_table(
        s,
        Cm(1.35),
        Cm(3.15),
        [Cm(5.0), Cm(8.0), Cm(8.0), Cm(8.5)],
        Cm(1.52),
        ["错误类型", "表象", "根因", "缓解策略"],
        [
            ("Wrong Column", "执行成功但指标不对", "字段同义、缩写、业务口径相近", "列注释、glossary、历史查询、人工标注链接"),
            ("Wrong Join", "结果重复或漏数", "join path 选择错误、桥表遗漏", "FK 图搜索、基数检查、结果 sanity check"),
            ("Wrong Value", "过滤条件找不到或误匹配", "值别名、大小写、脏数据", "值检索、模糊匹配、样例值提示"),
            ("Wrong Aggregation", "SUM/COUNT/AVG 搞错", "问题词到操作符链接不足", "IR sketch、操作符分类、单元测试"),
            ("Wrong Time Window", "日期范围偏移", "相对时间、财年、时区", "日期解析器、业务日历、执行前确认"),
        ],
        font_size=8.5,
    )
    add_card(s, Cm(1.45), Cm(12.35), Cm(30.5), Cm(3.2), "独立判断", "未来高质量 Text-to-SQL 系统的核心日志不应只记录 prompt 和 SQL，还应记录“链接证据链”：为什么选这些表列、排除了哪些候选、值是如何匹配的、执行验证通过了什么。", GREEN, PALE_GREEN, body_size=13.6)
    add_footer(s, page)


def production_arch(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "生产级架构：把 Schema Linking 做成可治理层", "落地方案")
    lanes = [
        ("数据资产层", "schema registry / lineage / glossary / permissions / sample values", BLUE),
        ("链接层", "retrieval / ranking / value linking / join graph / evidence store", GREEN),
        ("生成层", "planner / SQL generator / dialect adapter / constrained decoder", AMBER),
        ("验证层", "sandbox execution / cost guard / result tests / policy check / audit", RED),
    ]
    for i, (title, body, color) in enumerate(lanes):
        y = Cm(3.55 + i * 3.1)
        add_card(s, Cm(1.45), y, Cm(7.2), Cm(2.35), title, body, color, body_size=9.8, title_size=13.5)
        add_arrow(s, Cm(8.8), y + Cm(1.18), Cm(10.0), y + Cm(1.18), color=RGBColor(160, 174, 188), width=1.2)
        add_card(s, Cm(10.15), y, Cm(21.75), Cm(2.35), "关键产物", {
            0: "统一元数据、字段说明、数据血缘、权限标签、敏感字段标记、代表性值样本。",
            1: "候选 schema、链接分数、join path、值匹配证据、被排除候选及原因。",
            2: "SQL 草案、查询计划、方言转换、参数化过滤条件、候选集合。",
            3: "执行结果摘要、错误信息、空结果诊断、成本估计、审计日志和人类确认记录。",
        }[i], color, WHITE, body_size=10.8, title_size=12.6)
    add_text(s, Cm(1.5), Cm(16.15), Cm(29.8), Cm(0.8), "工程原则：不要把数据库语义永久藏在 prompt 里；它应该是可版本化、可测试、可审计的系统资产。", size=15.2, color=NAVY, bold=True)
    add_footer(s, page)


def design_principles(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "设计原则：让 LLM 少猜，让系统多给证据", "方法论")
    principles = [
        ("召回优先", "schema pruning 先保证关键表列不漏，再用排序控制噪声。", BLUE),
        ("证据可追踪", "每个 SQL 片段都能回溯到问题词、字段说明、样例值或业务规则。", GREEN),
        ("执行隔离", "只在沙箱或只读副本执行，设置超时、成本、行数与权限限制。", RED),
        ("多层验证", "语法、类型、执行、结果形状、业务单元测试分层检查。", AMBER),
        ("人机协同", "高风险查询让用户确认链接意图，而不是让模型自信地错。", PURPLE),
        ("持续学习", "把人工修正沉淀为 glossary、few-shot、单元测试和回归集。", CYAN),
    ]
    for i, (title, body, color) in enumerate(principles):
        row, col = divmod(i, 3)
        add_card(s, Cm(1.35 + col * 10.35), Cm(3.6 + row * 5.45), Cm(9.55), Cm(4.5), title, body, color, body_size=12.4)
    add_text(s, Cm(1.4), Cm(15.2), Cm(29.8), Cm(0.8), "一句话：高可靠 Text-to-SQL 不是让模型更“会说”，而是让它在更小、更准、更可验证的语义空间里行动。", size=15.2, color=NAVY, bold=True)
    add_footer(s, page)


def future(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "未来趋势：Text-to-SQL 会变成 Data Agent 的一部分", "趋势判断")
    trends = [
        ("从 SQL 到工作流", "真实问题常需要多步查询、临时表、数据清洗、可视化和报告生成。", BLUE),
        ("从 schema 到语义层", "指标口径、维度层级、业务规则会进入 semantic layer，LLM 只消费受控语义对象。", GREEN),
        ("从静态 prompt 到交互式检索", "模型可以主动请求更多元数据、样例值、文档和用户澄清。", CYAN),
        ("从准确率到治理指标", "可解释、合规、成本、延迟、审计、回归稳定性会与准确率同等重要。", AMBER),
        ("从单模型到组合系统", "小模型做检索/排序，大模型做推理，规则和执行器做约束。", PURPLE),
        ("从答案到信任", "系统要输出 SQL、结果，也要输出置信度、证据链和可复核路径。", RED),
    ]
    for i, (title, body, color) in enumerate(trends):
        row, col = divmod(i, 2)
        add_card(s, Cm(1.35 + col * 15.55), Cm(3.25 + row * 4.45), Cm(14.45), Cm(3.65), title, body, color, body_size=12.0)
    add_footer(s, page)


def summary(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "总结：Schema Linking 是 LLM Text-to-SQL 的可靠性中枢", "Takeaways")
    add_card(s, Cm(1.4), Cm(3.35), Cm(30.5), Cm(2.55), "1. 技术主线", "从规则映射、结构化解码、关系图编码，到 LLM + RAG + agent，演进方向始终围绕“如何把语言稳定绑定到数据结构”。", BLUE, PALE_BLUE, body_size=13.2)
    add_card(s, Cm(1.4), Cm(6.65), Cm(30.5), Cm(2.55), "2. 核心转折", "Spider 让跨 schema 泛化成为核心问题；BIRD 和 Spider 2.0 进一步暴露真实值、业务知识、超大 schema 与工作流复杂性。", GREEN, PALE_GREEN, body_size=13.2)
    add_card(s, Cm(1.4), Cm(9.95), Cm(30.5), Cm(2.55), "3. LLM 的正确用法", "不要把 LLM 当唯一解析器，而要把它放在可检索、可裁剪、可约束、可验证、可审计的系统中。", AMBER, PALE_AMBER, body_size=13.2)
    add_card(s, Cm(1.4), Cm(13.25), Cm(30.5), Cm(2.55), "4. 落地判断", "生产系统优先建设元数据治理、schema linking 证据链、执行沙箱和回归评测；模型升级只是其中一环。", RED, PALE_RED, body_size=13.2)
    add_footer(s, page)


def references(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "主要参考资料", "References")
    refs = [
        "Zhong et al., Seq2SQL: Generating Structured Queries from Natural Language using Reinforcement Learning, 2017.",
        "Xu et al., SQLNet: Generating Structured Queries From Natural Language Without Reinforcement Learning, arXiv:1711.04436.",
        "Yu et al., Spider: A Large-Scale Human-Labeled Dataset for Complex and Cross-Domain Semantic Parsing and Text-to-SQL Task, arXiv:1809.08887.",
        "Guo et al., Towards Complex Text-to-SQL in Cross-Domain Database with Intermediate Representation, arXiv:1905.08205.",
        "Wang et al., RAT-SQL: Relation-Aware Schema Encoding and Linking for Text-to-SQL Parsers, arXiv:1911.04942.",
        "Scholak et al., PICARD: Parsing Incrementally for Constrained Auto-Regressive Decoding from Language Models, arXiv:2109.05093.",
        "Li et al., RESDSQL: Decoupling Schema Linking and Skeleton Parsing for Text-to-SQL, arXiv:2302.05965.",
        "Pourreza & Rafiei, DIN-SQL: Decomposed In-Context Learning of Text-to-SQL with Self-Correction, arXiv:2304.11015.",
        "Li et al., BIRD: A Big Bench for Large-Scale Database Grounded Text-to-SQLs, arXiv:2305.03111.",
        "Gao et al., Text-to-SQL Empowered by Large Language Models: A Benchmark Evaluation / DAIL-SQL, arXiv:2308.15363.",
        "Wang et al., MAC-SQL: A Multi-Agent Collaborative Framework for Text-to-SQL, arXiv:2312.11242.",
        "Talaei et al., CHESS: Contextual Harnessing for Efficient SQL Synthesis, arXiv:2405.16755.",
        "Pourreza et al., CHASE-SQL: Multi-Path Reasoning and Preference Optimized Candidate Selection in Text-to-SQL, arXiv:2410.01943.",
        "Lei et al., Spider 2.0: Evaluating Language Models on Real-World Enterprise Text-to-SQL Workflows, arXiv:2411.07763.",
    ]
    add_bullets(s, Cm(1.35), Cm(3.15), Cm(30.8), Cm(12.8), refs, size=9.4, color=INK, gap=2)
    add_text(s, Cm(1.45), Cm(16.25), Cm(30), Cm(0.72), "注：本稿未直接复用论文图表，所有流程图和框架图均为基于文献整理后的原创示意。", size=10.5, color=SUBTLE)
    add_footer(s, page, "资料来源：上述公开论文、arXiv 摘要页和 benchmark 项目页。")


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_fns = [
        cover,
        thesis,
        agenda,
        evolution_timeline,
        task_definition,
        hard_cases,
        concept_mapping_linking,
        pre_neural,
        neural_wikisql,
        spider_shift,
        taxonomy,
        rat_sql,
        intermediate_ir,
        plm_resdsql,
        constrained_feedback,
        llm_paradigm,
        prompt_methods,
        rag_pruning,
        multi_agent,
        benchmarks,
        errors,
        production_arch,
        design_principles,
        future,
        summary,
        references,
    ]
    for idx, fn in enumerate(slide_fns, start=1):
        fn(prs, idx)
    prs.save(OUT)
    print(f"Generated {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
