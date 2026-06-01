from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


OUT = Path("/Users/wangqichen/Desktop/项目PPT_修改版.pptx")
ARCH_IMAGE = Path(__file__).resolve().parent / "assets" / "architecture-flow-image2.png"

W, H = Cm(33.867), Cm(19.05)

NAVY = RGBColor(18, 32, 47)
INK = RGBColor(32, 39, 47)
MUTED = RGBColor(91, 104, 118)
BLUE = RGBColor(32, 112, 196)
CYAN = RGBColor(29, 162, 184)
GREEN = RGBColor(51, 143, 99)
AMBER = RGBColor(218, 151, 45)
RED = RGBColor(190, 75, 75)
PALE = RGBColor(244, 248, 251)
LINE = RGBColor(209, 220, 229)
WHITE = RGBColor(255, 255, 255)

FONT = "PingFang SC"

COMMON_SOURCE = (
    "资料来源：项目申请书；市场数据核查见《备用.docx》：机器人大讲堂(2024-10-14)、搜狐/智研咨询、"
    "WiseGuy Reports、东方财富网/新浪财经(2025-10-31)、深圳市人工智能产业协会官方平台。"
)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, x, y, w, h, text, size=20, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.16)
    tf.margin_right = Cm(0.16)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.line_spacing = 1.05
    if align:
        p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_title(slide, title, kicker=None):
    if kicker:
        add_text(slide, Cm(1.25), Cm(0.72), Cm(12), Cm(0.52), kicker, size=11, color=BLUE, bold=True)
    add_text(slide, Cm(1.2), Cm(1.2), Cm(25), Cm(1.05), title, size=27, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.2), Cm(2.44), Cm(2.15), Cm(0.08))
    set_fill(line, CYAN)
    line.line.fill.background()


def add_footer(slide, extra=None):
    footer = COMMON_SOURCE if extra is None else extra
    add_text(slide, Cm(1.05), Cm(18.12), Cm(31.8), Cm(0.55), footer, size=6.4, color=RGBColor(105, 116, 127))
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.05), Cm(17.94), Cm(31.8), Cm(0.02))
    set_fill(rule, LINE)
    rule.line.fill.background()


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 252, 254)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, body_size=11.6):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    set_fill(card, WHITE)
    set_line(card)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Cm(0.12), h)
    set_fill(bar, accent)
    bar.line.fill.background()
    add_text(slide, x + Cm(0.42), y + Cm(0.35), w - Cm(0.7), Cm(0.62), title, size=16, color=NAVY, bold=True)
    add_text(slide, x + Cm(0.42), y + Cm(1.15), w - Cm(0.7), h - Cm(1.3), body, size=body_size, color=MUTED)
    return card


def add_metric(slide, x, y, w, h, value, label, color=BLUE):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.07
    set_fill(box, RGBColor(248, 251, 253))
    set_line(box, RGBColor(220, 228, 235))
    add_text(slide, x + Cm(0.25), y + Cm(0.34), w - Cm(0.5), Cm(0.88), value, size=25, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Cm(0.3), y + Cm(1.36), w - Cm(0.6), Cm(0.7), label, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_pill(slide, x, y, w, text, color):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Cm(0.72))
    pill.adjustments[0] = 0.35
    set_fill(pill, color)
    pill.line.fill.background()
    add_text(slide, x + Cm(0.12), y + Cm(0.13), w - Cm(0.24), Cm(0.4), text, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(136, 154, 171), width=1.6):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_line(slide, x1, y1, x2, y2, color=RGBColor(136, 154, 171), width=1.6):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_bullets(slide, x, y, w, h, bullets, size=14, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(4)
        p.line_spacing = 1.05
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(8.4), H)
    set_fill(left, NAVY)
    left.line.fill.background()
    for i in range(7):
        line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.2), Cm(4.0 + i * 1.15), Cm(5.7), Cm(0.025))
        set_fill(line, RGBColor(54, 82, 108))
        line.line.fill.background()
    add_text(s, Cm(1.3), Cm(1.15), Cm(4.5), Cm(0.55), "创新大赛校赛", size=13, color=RGBColor(196, 221, 236), bold=True)
    add_text(s, Cm(10.6), Cm(4.9), Cm(19.5), Cm(2.25), "复杂环境中\n多机器人运动规划与探索研究", size=34, color=NAVY, bold=True)
    add_text(s, Cm(10.75), Cm(9.2), Cm(18.5), Cm(0.85), "团队：张婧瑶、陈可欣、闫天路、张释文、荣祎铭", size=16, color=INK, bold=True)
    add_text(s, Cm(10.75), Cm(10.42), Cm(10.0), Cm(0.75), "指导老师：待补充", size=15, color=MUTED)
    add_text(s, Cm(10.75), Cm(12.0), Cm(10.0), Cm(0.65), "同济大学本科生项目组", size=13.5, color=MUTED)
    add_footer(s, "资料来源：项目申请书。")


def slide_applications(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "应用场景：空地协同解决真实作业需求", "应用场景")
    add_text(s, Cm(1.25), Cm(3.08), Cm(29.5), Cm(0.75), "本页只回答“为什么值得做”：空中全局视野与地面精细执行互补，适合高风险、高复杂度、高频巡检任务。", size=15.2, color=MUTED)
    scenarios = [
        ("灾后搜救与应急救援", "地震、矿难等场景环境破碎、信号微弱、人工进入风险高。无人机快速建图与搜索，地面机器人深入危险区域执行探测与物资投放。", RED),
        ("大型工业园区智慧巡检", "化工厂、仓储物流园区、油气管网需要高频巡检。无人机快速筛查异常点，UGV 进行近距离复检和数据采集。", BLUE),
        ("城市治理与安防监控", "大型活动安保、边防巡逻、交通管理需要全域监测。空中宏观预警与地面微观精检形成联动响应。", GREEN),
    ]
    for i, (title, body, color) in enumerate(scenarios):
        add_card(s, Cm(1.35 + i * 10.45), Cm(4.55), Cm(9.45), Cm(7.3), title, body, accent=color, body_size=12.2)
    add_text(s, Cm(2.1), Cm(13.25), Cm(29), Cm(0.9), "核心价值：更快覆盖、更少人工暴露、更强复杂环境适应能力。", size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_challenges(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "技术挑战：空地异构协同为什么难", "问题定义")
    data = [
        ("平台异构", "UAV 与 UGV 在感知范围、运动约束、任务能力上差异显著，需要统一环境表达与协同决策接口。", BLUE),
        ("动态环境", "障碍、人员和通道状态持续变化，规划必须跟随感知实时更新，避免信息滞后。", RED),
        ("连续空间冲突", "多机器人在狭窄或高密度场景中易发生碰撞、死锁与执行偏差，需要显式冲突消解。", AMBER),
        ("弱通信与安全", "不能依赖高频稳定通信；学习策略可提升效率，但输出必须被规划控制层约束成安全轨迹。", GREEN),
    ]
    for idx, (title, body, color) in enumerate(data):
        row, col = divmod(idx, 2)
        add_card(s, Cm(1.5 + col * 15.5), Cm(3.75 + row * 5.35), Cm(14.2), Cm(4.55), title, body, accent=color, body_size=12.3)
    center = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(10.8), Cm(8.15), Cm(12.0), Cm(1.95))
    center.adjustments[0] = 0.12
    set_fill(center, NAVY)
    center.line.fill.background()
    add_text(s, Cm(11.1), Cm(8.52), Cm(11.4), Cm(0.76), "目标：未知动态环境下的安全协同探索", size=15.3, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def flow_box(slide, x, y, w, h, title, body, color):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    set_fill(box, WHITE)
    set_line(box)
    add_pill(slide, x + Cm(0.35), y + Cm(0.32), Cm(3.0), title, color)
    add_text(slide, x + Cm(0.35), y + Cm(1.22), w - Cm(0.7), h - Cm(1.35), body, size=12.2, color=INK, bold=True, align=PP_ALIGN.CENTER)
    return box


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    s.shapes.add_picture(str(ARCH_IMAGE), Cm(0), Cm(0), width=W, height=H)
    title_band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), W, Cm(2.85))
    set_fill(title_band, WHITE, transparency=9)
    title_band.line.fill.background()
    add_title(s, "技术架构：从感知到执行的闭环流程", "技术方案")

    labels = [
        (Cm(5.0), Cm(8.52), Cm(3.35), Cm(0.92), "多源感知\n点云 / 位姿 / 风险", BLUE),
        (Cm(10.65), Cm(8.52), Cm(3.75), Cm(0.92), "统一环境模型\n几何 / 语义 / 风险", CYAN),
        (Cm(16.05), Cm(8.52), Cm(4.1), Cm(0.92), "探索与任务分配\n高信息区域评估", GREEN),
        (Cm(21.55), Cm(8.52), Cm(3.95), Cm(0.92), "全局冲突消解\nCBS / K-CBS", AMBER),
        (Cm(27.05), Cm(8.52), Cm(3.85), Cm(0.92), "局部轨迹优化\nMPC 滚动规划", BLUE),
        (Cm(14.75), Cm(14.05), Cm(5.95), Cm(0.92), "空地执行与验证\nUAV 侦察 + UGV 作业", NAVY),
    ]
    for x, y, w, h, txt, color in labels:
        add_text(s, x, y, w, h, txt, size=11.6, color=color, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, Cm(3.1), Cm(14.9), Cm(5.6), Cm(0.48), "执行反馈更新地图", size=10.6, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Cm(22.8), Cm(13.78), Cm(6.3), Cm(0.48), "约束下生成可执行轨迹", size=10.6, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def draw_drone(slide, x, y, color, label):
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Cm(0.48), Cm(0.48))
    set_fill(body, color)
    body.line.fill.background()
    add_line(slide, x - Cm(0.28), y + Cm(0.24), x + Cm(0.76), y + Cm(0.24), color=color, width=1.6)
    add_line(slide, x + Cm(0.24), y - Cm(0.28), x + Cm(0.24), y + Cm(0.76), color=color, width=1.6)
    add_text(slide, x - Cm(0.45), y + Cm(0.62), Cm(1.4), Cm(0.4), label, size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)


def slide_cbmpc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "CB-MPC：全局冲突树搜索 + 局部滚动优化", "核心算法")
    add_text(
        s, Cm(1.25), Cm(3.05), Cm(30.5), Cm(0.8),
        "CB-MPC = Conflict-Based Search（全局冲突树搜索） + Model Predictive Control（局部滚动优化）。",
        size=17, color=NAVY, bold=True,
    )
    add_card(
        s, Cm(1.35), Cm(4.35), Cm(14.35), Cm(3.3),
        "CBS 做什么",
        "检测多机器人之间的路径冲突，并在冲突树中生成时间、空间或机器人级约束，给后续重规划划定安全边界。",
        accent=RED, body_size=12.2,
    )
    add_card(
        s, Cm(17.0), Cm(4.35), Cm(14.35), Cm(3.3),
        "MPC 做什么",
        "在 CBS 给出的约束下，为每台机器人滚动生成满足动力学约束、控制约束和安全距离要求的平滑轨迹。",
        accent=GREEN, body_size=12.2,
    )

    # Diagram stage 1: conflict detection.
    stage_y = Cm(9.0)
    stage_w = Cm(8.9)
    for i, (title, color) in enumerate([("1 检测到路径冲突", RED), ("2 CBS 生成约束", AMBER), ("3 MPC 各自重规划绕开", GREEN)]):
        x = Cm(1.35 + i * 10.45)
        box = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, stage_y, stage_w, Cm(5.6))
        box.adjustments[0] = 0.06
        set_fill(box, WHITE)
        set_line(box)
        add_pill(s, x + Cm(0.45), stage_y + Cm(0.38), Cm(4.4), title, color)

    # Stage 1 paths.
    x = Cm(1.35)
    add_line(s, x + Cm(1.2), stage_y + Cm(4.55), x + Cm(7.55), stage_y + Cm(1.35), BLUE, 2.1)
    add_line(s, x + Cm(1.2), stage_y + Cm(1.35), x + Cm(7.55), stage_y + Cm(4.55), RED, 2.1)
    conflict = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Cm(4.12), stage_y + Cm(2.65), Cm(0.72), Cm(0.72))
    set_fill(conflict, RED, transparency=10)
    set_line(conflict, RED, 2)
    add_text(s, x + Cm(3.05), stage_y + Cm(3.45), Cm(2.95), Cm(0.45), "同一时空占用", size=9.8, color=RED, bold=True, align=PP_ALIGN.CENTER)
    draw_drone(s, x + Cm(0.78), stage_y + Cm(4.25), BLUE, "UAV A")
    draw_drone(s, x + Cm(0.78), stage_y + Cm(1.05), RED, "UAV B")

    # Stage 2 constraints.
    x = Cm(11.8)
    tree = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Cm(1.1), stage_y + Cm(1.35), Cm(6.4), Cm(0.85))
    tree.adjustments[0] = 0.08
    set_fill(tree, RGBColor(255, 248, 235))
    set_line(tree, AMBER)
    add_text(s, x + Cm(1.25), stage_y + Cm(1.54), Cm(6.1), Cm(0.38), "冲突树节点：A 与 B 冲突", size=10.3, color=INK, bold=True, align=PP_ALIGN.CENTER)
    for j, txt in enumerate(["约束 A：t 时刻避开冲突区", "约束 B：t 时刻避开冲突区"]):
        c = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Cm(0.8 + j * 4.0), stage_y + Cm(3.4), Cm(3.25), Cm(0.95))
        c.adjustments[0] = 0.08
        set_fill(c, RGBColor(252, 249, 242))
        set_line(c, AMBER)
        add_text(s, x + Cm(0.9 + j * 4.0), stage_y + Cm(3.58), Cm(3.05), Cm(0.45), txt, size=8.8, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_arrow(s, x + Cm(4.3), stage_y + Cm(2.2), x + Cm(2.4 + j * 4.0), stage_y + Cm(3.4), AMBER, 1.2)

    # Stage 3 replanning.
    x = Cm(22.25)
    add_line(s, x + Cm(1.0), stage_y + Cm(4.6), x + Cm(3.4), stage_y + Cm(3.65), BLUE, 2.1)
    add_line(s, x + Cm(3.4), stage_y + Cm(3.65), x + Cm(5.3), stage_y + Cm(2.0), BLUE, 2.1)
    add_line(s, x + Cm(5.3), stage_y + Cm(2.0), x + Cm(7.6), stage_y + Cm(1.25), BLUE, 2.1)
    add_line(s, x + Cm(1.0), stage_y + Cm(1.25), x + Cm(3.4), stage_y + Cm(2.25), RED, 2.1)
    add_line(s, x + Cm(3.4), stage_y + Cm(2.25), x + Cm(5.3), stage_y + Cm(3.9), RED, 2.1)
    add_line(s, x + Cm(5.3), stage_y + Cm(3.9), x + Cm(7.6), stage_y + Cm(4.6), RED, 2.1)
    safe = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Cm(4.0), stage_y + Cm(2.65), Cm(0.85), Cm(0.85))
    set_fill(safe, RGBColor(237, 248, 241))
    set_line(safe, GREEN, 1.5)
    add_text(s, x + Cm(3.0), stage_y + Cm(3.48), Cm(2.9), Cm(0.45), "绕开冲突区", size=9.8, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    draw_drone(s, x + Cm(0.65), stage_y + Cm(4.3), BLUE, "UAV A")
    draw_drone(s, x + Cm(0.65), stage_y + Cm(0.98), RED, "UAV B")

    add_arrow(s, Cm(10.25), stage_y + Cm(2.85), Cm(11.45), stage_y + Cm(2.85), width=1.8)
    add_arrow(s, Cm(20.7), stage_y + Cm(2.85), Cm(21.9), stage_y + Cm(2.85), width=1.8)
    add_footer(s)


def slide_algorithm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "核心技术路径", "算法体系")
    add_text(s, Cm(1.25), Cm(3.1), Cm(18), Cm(0.8), "从实时建图到轨迹执行，系统把全局协同与局部控制放在同一条约束链上。", size=15.5, color=MUTED)
    rows = [
        ("FAST-LIVO2 实时三维建图", "融合 LiDAR / Visual / IMU 信息，输出动态地图；当前阶段使用仿真点云验证规划算法，后续接入实时建图。", BLUE),
        ("MPC 局部轨迹生成", "面向异构机器人生成满足动力学约束、控制约束与安全距离要求的平滑轨迹。", GREEN),
        ("K-CBS 思想 + CB-MPC 框架", "兼顾全局冲突搜索与局部滚动优化，处理连续空间中的碰撞、死锁与执行偏差。", AMBER),
        ("RACER 分布式探索机制", "引入去中心化探索思路，提高扩展性，降低单点故障对全局任务的影响。", CYAN),
    ]
    for i, (title, body, color) in enumerate(rows):
        add_card(s, Cm(1.35), Cm(4.35 + i * 2.85), Cm(24.2), Cm(2.2), title, body, accent=color)
        if i < 3:
            add_arrow(s, Cm(26.0), Cm(5.45 + i * 2.85), Cm(26.0), Cm(6.05 + i * 2.85), color)
    add_metric(s, Cm(27.0), Cm(5.0), Cm(4.7), Cm(2.25), "安全", "显式约束", RED)
    add_metric(s, Cm(27.0), Cm(7.75), Cm(4.7), Cm(2.25), "实时", "滚动规划", BLUE)
    add_metric(s, Cm(27.0), Cm(10.5), Cm(4.7), Cm(2.25), "可扩展", "分布式协同", GREEN)
    add_footer(s)


def slide_innovation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "两项核心创新", "创新点")
    add_card(s, Cm(1.45), Cm(3.75), Cm(14.85), Cm(10.4), "1. 实时“感知-规划”紧耦合", "FAST-LIVO2 动态地图输出直接馈入 MPC 规划器；规划指令与感知数据闭环更新；基于动态障碍物短期预测主动调整局部路径。\n\n价值：从“被动避障”走向“前瞻性避障”。", BLUE, body_size=14.2)
    add_card(s, Cm(17.1), Cm(3.75), Cm(14.85), Cm(10.4), "2. 异构协同冲突解决范式", "CB-MPC 兼顾全局路径最优性与局部实时调整；把异构平台运动学差异纳入冲突检测与消解；后续接入 RACER 后可扩展至去中心化集群探索。\n\n价值：支持 UAV / UGV 精细化轨迹协同。", GREEN, body_size=14.2)
    add_footer(s)


def slide_progress(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "项目进展：仿真验证形成直观结果", "项目进展")
    add_text(s, Cm(1.45), Cm(3.12), Cm(30.6), Cm(1.1), "80个随机障碍物 · 4架无人机 · 全程无碰撞", size=29, color=RED, bold=True, align=PP_ALIGN.CENTER)

    ph = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.55), Cm(5.0), Cm(18.8), Cm(9.6))
    ph.adjustments[0] = 0.04
    set_fill(ph, RGBColor(246, 249, 252))
    set_line(ph, RGBColor(154, 172, 190), 1.4)
    ph.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_text(s, Cm(1.85), Cm(8.75), Cm(18.2), Cm(0.8), "RViz 截图占位", size=23, color=RGBColor(126, 143, 160), bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Cm(1.85), Cm(9.75), Cm(18.2), Cm(0.55), "视频见附件", size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

    add_card(s, Cm(21.25), Cm(5.0), Cm(10.45), Cm(2.35), "核心实现", "CB-MPC 多机器人协同运动规划算法已完成核心实现与仿真验证。", BLUE, body_size=11.2)
    add_card(s, Cm(21.25), Cm(7.85), Cm(10.45), Cm(2.35), "实时指标", "规划频率 2Hz，控制频率 30Hz，满足当前仿真场景实时性要求。", GREEN, body_size=11.2)
    add_card(s, Cm(21.25), Cm(10.7), Cm(10.45), Cm(2.35), "系统对接", "已完成 CB-MPC 与 EGO-Planner、ROS 话题通信、轨迹服务器和 RViz 可视化集成。", AMBER, body_size=11.2)

    add_metric(s, Cm(2.0), Cm(15.25), Cm(5.2), Cm(1.9), "4架", "无人机协同", BLUE)
    add_metric(s, Cm(7.7), Cm(15.25), Cm(5.2), Cm(1.9), "80个", "随机障碍物", AMBER)
    add_metric(s, Cm(13.4), Cm(15.25), Cm(5.2), Cm(1.9), "0碰撞", "安全到达目标", GREEN)
    add_footer(s)


def slide_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "市场前景：低空经济与机器人产业交汇", "市场数据")
    add_text(s, Cm(1.25), Cm(3.1), Cm(29.5), Cm(0.8), "备用文档显示，部分原始市场预测口径需谨慎引用；本页仅保留可查证或有多源印证的数据。", size=14.5, color=MUTED)
    add_text(s, Cm(1.35), Cm(4.45), Cm(12.5), Cm(0.65), "中国智能巡检机器人市场规模", size=17, color=NAVY, bold=True)
    max_w = Cm(10.2)
    for i, (year, val, color) in enumerate([(2023, 19.71, BLUE), (2025, 27.48, GREEN)]):
        y = Cm(5.45 + i * 1.25)
        add_text(s, Cm(1.6), y + Cm(0.08), Cm(1.45), Cm(0.45), str(year), size=11.5, color=MUTED)
        bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(3.15), y, max_w, Cm(0.48))
        set_fill(bg, RGBColor(232, 238, 244))
        bg.line.fill.background()
        bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(3.15), y, max_w * (val / 30), Cm(0.48))
        set_fill(bar, color)
        bar.line.fill.background()
        add_text(s, Cm(3.35) + max_w * (val / 30), y - Cm(0.02), Cm(2.2), Cm(0.55), f"{val:.2f}亿元", size=12.5, color=INK, bold=True)
    add_text(s, Cm(1.6), Cm(8.15), Cm(12.5), Cm(0.8), "2025 年采用智研咨询/搜狐转载口径；备用文档提示“32亿元”未获直接印证。", size=9.5, color=MUTED)
    add_text(s, Cm(17.0), Cm(4.45), Cm(12.5), Cm(0.65), "全球工业级无人机市场", size=17, color=NAVY, bold=True)
    for i, (year, val, color) in enumerate([(2025, 87.9, CYAN), (2035, 300, AMBER)]):
        y = Cm(5.45 + i * 1.25)
        add_text(s, Cm(17.25), y + Cm(0.08), Cm(1.45), Cm(0.45), str(year), size=11.5, color=MUTED)
        bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(18.8), y, max_w, Cm(0.48))
        set_fill(bg, RGBColor(232, 238, 244))
        bg.line.fill.background()
        bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(18.8), y, max_w * (val / 320), Cm(0.48))
        set_fill(bar, color)
        bar.line.fill.background()
        add_text(s, Cm(19.0) + max_w * (val / 320), y - Cm(0.02), Cm(3.1), Cm(0.55), f"{val:g}亿美元", size=12.5, color=INK, bold=True)
    add_text(s, Cm(17.25), Cm(8.15), Cm(12.5), Cm(0.8), "WiseGuy Reports 预测：2025 年 87.9 亿美元，2035 年 300 亿美元。", size=9.5, color=MUTED)
    add_card(s, Cm(1.35), Cm(10.25), Cm(14.7), Cm(3.5), "产业融合窗口", "空地一体集群智慧 AI 场景解决方案融合低空经济与机器人两大万亿级产业。", BLUE)
    add_card(s, Cm(17.0), Cm(10.25), Cm(14.7), Cm(3.5), "应用市场空间", "范丛明公开发言：相关场景解决方案未来 5 年将催生超百亿元规模的应用市场。", GREEN)
    add_footer(s, "数据来源：机器人大讲堂(2024-10-14)；搜狐网转载智研咨询报告；WiseGuy Reports《Industrial Grade Drone Market Insights》；东方财富网/新浪财经(2025-10-31)；深圳市人工智能产业协会官方平台。")


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "团队分工：建图组 × 规划组", "团队介绍")
    add_card(s, Cm(1.35), Cm(3.7), Cm(14.7), Cm(6.6), "建图组", "张婧瑶：项目统筹、仿真环境搭建、ROS 系统集成、CB-MPC 与 EGO-Planner 对接\n张释文：建图模块数据处理与环境建模", BLUE, body_size=13)
    add_card(s, Cm(17.0), Cm(3.7), Cm(14.7), Cm(6.6), "规划组", "闫天路：CB-MPC 核心算法实现与调试\n陈可欣：规划算法设计与文献调研\n荣祎铭：规划算法搭建与仿真验证", GREEN, body_size=13)
    add_text(s, Cm(2.0), Cm(12.0), Cm(29.0), Cm(1.0), "团队能力覆盖 ROS、运动规划、三维感知、环境建模与仿真验证，已形成算法研究和系统集成协作闭环。", size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "后续规划与预期成果", "实施路径")
    steps = [
        ("2026.05\n已完成", "CB-MPC 核心实现\n4机仿真验证\nEGO-Planner 对接", BLUE),
        ("2026.06-08", "接入 FAST-LIVO2\n加入 UGV 模型\n动态障碍验证", GREEN),
        ("2026.09-11", "强化学习探索策略\n大规模仿真对比\n真实平台初测", AMBER),
        ("2026.11-2027.03", "室内/户外实物实验\n系统迭代优化\n成果完善", RED),
    ]
    for i, (time, body, color) in enumerate(steps):
        x = Cm(1.45 + i * 7.8)
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Cm(2.15), Cm(4.0), Cm(1.1), Cm(1.1))
        set_fill(circ, color)
        circ.line.fill.background()
        if i < 3:
            line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Cm(3.25), Cm(4.52), Cm(5.65), Cm(0.08))
            set_fill(line, RGBColor(205, 217, 229))
            line.line.fill.background()
        add_text(s, x, Cm(5.55), Cm(5.6), Cm(1.25), time, size=15.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(7.25), Cm(5.6), Cm(2.6), body, size=12.2, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(s, Cm(1.5), Cm(12.0), Cm(14.7), Cm(3.3), "算法成果", "形成面向空地异构多机器人系统的统一环境建模、协同决策、冲突处理与自主探索方法。", BLUE)
    add_card(s, Cm(17.05), Cm(12.0), Cm(14.7), Cm(3.3), "系统成果", "完成仿真与典型真实场景实验验证，为巡检、环境探索等应用提供技术参考。", GREEN)
    add_footer(s)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    block = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), H)
    set_fill(block, NAVY)
    block.line.fill.background()
    add_text(s, Cm(2.2), Cm(6.65), Cm(29.5), Cm(2.6), "敬请各位老师批评指正！", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Cm(2.2), Cm(10.0), Cm(29.5), Cm(0.9), "复杂环境中多机器人运动规划与探索研究", size=17, color=RGBColor(202, 221, 236), align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_cover(prs)
    slide_applications(prs)
    slide_challenges(prs)
    slide_architecture(prs)
    slide_cbmpc(prs)
    slide_algorithm(prs)
    slide_innovation(prs)
    slide_progress(prs)
    slide_market(prs)
    slide_team(prs)
    slide_roadmap(prs)
    slide_thanks(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
