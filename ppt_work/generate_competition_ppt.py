from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


OUT = Path("/Users/wangqichen/Desktop/空地异构多机器人协同系统_项目PPT.pptx")

W, H = Cm(33.867), Cm(19.05)

NAVY = RGBColor(18, 32, 47)
INK = RGBColor(32, 39, 47)
MUTED = RGBColor(91, 104, 118)
BLUE = RGBColor(32, 112, 196)
CYAN = RGBColor(29, 162, 184)
GREEN = RGBColor(51, 143, 99)
AMBER = RGBColor(218, 151, 45)
RED = RGBColor(190, 75, 75)
PALE = RGBColor(244, 248, 251)
LINE = RGBColor(209, 220, 229)
WHITE = RGBColor(255, 255, 255)

FONT = "PingFang SC"
TITLE_FONT = "PingFang SC"

COMMON_SOURCE = (
    "资料来源：项目申请书；市场数据核查见《备用.docx》：机器人大讲堂(2024-10-14)、搜狐/智研咨询、"
    "WiseGuy Reports、东方财富网/新浪财经(2025-10-31)、深圳市人工智能产业协会官方平台。"
)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def text_frame(shape, text, size=20, color=INK, bold=False, align=None, line_spacing=1.05):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Cm(0.18)
    tf.margin_right = Cm(0.18)
    tf.margin_top = Cm(0.08)
    tf.margin_bottom = Cm(0.08)
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tf


def add_text(slide, x, y, w, h, text, size=20, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    text_frame(box, text, size=size, color=color, bold=bold, align=align)
    return box


def add_title(slide, title, kicker=None):
    if kicker:
        add_text(slide, Cm(1.25), Cm(0.72), Cm(12), Cm(0.52), kicker, size=11, color=BLUE, bold=True)
    add_text(slide, Cm(1.2), Cm(1.2), Cm(24), Cm(1.05), title, size=27, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.2), Cm(2.44), Cm(2.15), Cm(0.08))
    set_fill(line, CYAN)
    line.line.fill.background()


def add_footer(slide, extra=None):
    footer = COMMON_SOURCE if extra is None else extra
    add_text(slide, Cm(1.05), Cm(18.12), Cm(31.8), Cm(0.55), footer, size=6.4, color=RGBColor(105, 116, 127))
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.05), Cm(17.94), Cm(31.8), Cm(0.02))
    set_fill(rule, LINE)
    rule.line.fill.background()


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 252, 254)


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    set_fill(card, WHITE)
    set_line(card)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Cm(0.12), h)
    set_fill(bar, accent)
    bar.line.fill.background()
    add_text(slide, x + Cm(0.42), y + Cm(0.35), w - Cm(0.7), Cm(0.55), title, size=16, color=NAVY, bold=True)
    add_text(slide, x + Cm(0.42), y + Cm(1.13), w - Cm(0.7), h - Cm(1.32), body, size=11.6, color=MUTED)
    return card


def add_metric(slide, x, y, w, h, value, label, color=BLUE):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.07
    set_fill(box, RGBColor(248, 251, 253))
    set_line(box, RGBColor(220, 228, 235))
    add_text(slide, x + Cm(0.25), y + Cm(0.35), w - Cm(0.5), Cm(0.85), value, size=25, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Cm(0.3), y + Cm(1.36), w - Cm(0.6), Cm(0.7), label, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_pill(slide, x, y, w, text, color):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Cm(0.72))
    pill.adjustments[0] = 0.35
    set_fill(pill, color)
    pill.line.fill.background()
    add_text(slide, x + Cm(0.12), y + Cm(0.13), w - Cm(0.24), Cm(0.4), text, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(136, 154, 171)):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True
    return line


def add_bullets(slide, x, y, w, h, bullets, size=14, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(4)
        p.line_spacing = 1.05
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(11.1), H)
    set_fill(left, NAVY)
    left.line.fill.background()
    # Abstract exploration grid.
    for i in range(8):
        for j in range(6):
            r = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(2.0 + i * 0.95), Cm(3.0 + j * 0.95), Cm(0.1), Cm(0.1))
            set_fill(r, RGBColor(61, 91, 119))
            r.line.fill.background()
    add_pill(s, Cm(1.55), Cm(1.22), Cm(4.2), "创新大赛校赛", CYAN)
    add_text(s, Cm(1.45), Cm(3.2), Cm(8.6), Cm(4.9), "空地异构\n多机器人协同系统", size=34, color=WHITE, bold=True)
    add_text(s, Cm(1.52), Cm(8.5), Cm(7.8), Cm(1.2), "面向未知动态环境的协同运动规划与自主探索", size=15, color=RGBColor(204, 221, 235))

    add_text(s, Cm(12.4), Cm(3.05), Cm(18.5), Cm(1.05), "空中侦察 + 地面执行", size=31, color=NAVY, bold=True)
    add_text(s, Cm(12.45), Cm(4.35), Cm(15.4), Cm(0.85), "以 UAV 与 UGV 为核心载体，构建“感知-规划-协同”一体化技术架构", size=15.2, color=MUTED)
    add_metric(s, Cm(12.45), Cm(6.05), Cm(5.35), Cm(2.45), "2Hz", "滚动规划频率", BLUE)
    add_metric(s, Cm(18.2), Cm(6.05), Cm(5.35), Cm(2.45), "30Hz", "控制执行频率", GREEN)
    add_metric(s, Cm(23.95), Cm(6.05), Cm(5.35), Cm(2.45), "4机", "仿真避碰验证", AMBER)

    for idx, (label, detail, color) in enumerate([
        ("实时三维建图", "FAST-LIVO2 接入规划闭环", BLUE),
        ("安全协同控制", "CB-MPC 处理连续空间冲突", GREEN),
        ("分布式探索", "RACER 思路提升扩展性", AMBER),
    ]):
        add_card(s, Cm(12.45 + idx * 5.85), Cm(10.05), Cm(5.3), Cm(2.55), label, detail, accent=color)
    add_text(s, Cm(12.45), Cm(14.2), Cm(17), Cm(0.7), "团队：同济大学本科生项目组｜建图组 × 规划组", size=14.5, color=NAVY, bold=True)
    add_footer(s)


def slide_context(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "为什么需要空地异构协同", "项目背景")
    add_text(s, Cm(1.25), Cm(3.2), Cm(30.5), Cm(0.95), "单一机器人在复杂场景中容易被视野、通行能力、感知精度和任务负载限制；空地协同则把全局视野与近场执行能力组合起来。", size=16.5, color=MUTED)
    items = [
        ("灾后搜救", "环境破碎、信号微弱、人工进入风险高；无人机先行建图，地面机器人深入探测与投放。", RED),
        ("工业巡检", "化工、仓储、油气管网需要高频巡检；空中快速筛查，地面复检异常点。", BLUE),
        ("城市治理", "大型活动安保、边防巡逻、交通管理需要宏观预警与微观精检联动。", GREEN),
    ]
    for i, (title, body, color) in enumerate(items):
        add_card(s, Cm(1.35 + i * 10.4), Cm(5.0), Cm(9.45), Cm(5.0), title, body, accent=color)
    add_text(s, Cm(1.25), Cm(11.2), Cm(8), Cm(0.8), "行业共性痛点", size=18, color=NAVY, bold=True)
    add_bullets(s, Cm(1.35), Cm(12.25), Cm(30.5), Cm(3.6), [
        "未知动态环境：障碍、人员、通道状态持续变化，规划必须实时更新。",
        "信息不完全与弱通信：探索过程中无法依赖持续稳定的集中式调度。",
        "平台异构：UAV 与 UGV 的感知范围、动力学约束、任务能力差异显著。",
        "安全可执行性：学习方法可提升效率，但必须被规划与控制层约束住。"
    ], size=15)
    add_footer(s)


def slide_painpoints(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "项目聚焦的四个关键问题", "问题定义")
    data = [
        ("统一环境建模", "如何让 UAV 的空中快速探索与 UGV 的地面安全通行共享同一套环境表达？", BLUE),
        ("连续空间冲突", "狭窄或高密度场景下，如何避免碰撞、死锁与执行偏差？", RED),
        ("动态任务分配", "探索区域分配不均、重复覆盖、频繁重规划如何被实时抑制？", AMBER),
        ("学习结果安全化", "强化学习能提升探索效率，但如何保证输出轨迹可执行、可验证？", GREEN),
    ]
    for idx, (title, body, color) in enumerate(data):
        row, col = divmod(idx, 2)
        add_card(s, Cm(1.5 + col * 15.5), Cm(3.7 + row * 5.55), Cm(14.2), Cm(4.65), title, body, accent=color)
    center = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(11.75), Cm(8.25), Cm(10.2), Cm(2.0))
    center.adjustments[0] = 0.12
    set_fill(center, NAVY)
    center.line.fill.background()
    add_text(s, Cm(12.05), Cm(8.62), Cm(9.6), Cm(0.8), "目标：未知动态环境下的安全协同探索", size=15.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "“感知-规划-协同”三位一体架构", "技术方案")
    blocks = [
        ("多源感知", "UAV / UGV\n点云、位姿、风险", BLUE),
        ("统一环境模型", "几何结构\n语义与风险表征", CYAN),
        ("探索与分配", "高信息区域评估\n动态任务分配", AMBER),
        ("协同规划控制", "K-CBS / CB-MPC\n动力学与安全约束", GREEN),
        ("空地执行", "空中侦察\n地面精细作业", RED),
    ]
    y = Cm(5.15)
    for i, (title, body, color) in enumerate(blocks):
        x = Cm(1.2 + i * 6.35)
        card = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Cm(5.1), Cm(4.25))
        card.adjustments[0] = 0.08
        set_fill(card, WHITE)
        set_line(card)
        add_pill(s, x + Cm(0.5), y + Cm(0.45), Cm(2.8), title, color)
        add_text(s, x + Cm(0.45), y + Cm(1.55), Cm(4.2), Cm(1.55), body, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
        if i < len(blocks) - 1:
            add_arrow(s, x + Cm(5.25), y + Cm(2.1), x + Cm(6.05), y + Cm(2.1))
    add_card(s, Cm(2.0), Cm(11.2), Cm(9.2), Cm(3.2), "闭环反馈", "地图动态更新直接影响局部轨迹与任务分配，避免传统串行流水线的信息滞后。", BLUE)
    add_card(s, Cm(12.25), Cm(11.2), Cm(9.2), Cm(3.2), "弱通信适配", "通过分布式探索机制降低对集中调度和高频通信的依赖。", GREEN)
    add_card(s, Cm(22.5), Cm(11.2), Cm(9.2), Cm(3.2), "异构约束显式化", "将 UAV / UGV 的运动学差异纳入冲突检测与消解模型。", AMBER)
    add_footer(s)


def slide_algorithm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "核心技术路径", "算法体系")
    add_text(s, Cm(1.25), Cm(3.1), Cm(18), Cm(0.8), "从实时建图到轨迹执行，系统把全局协同与局部控制放在同一条约束链上。", size=15.5, color=MUTED)
    rows = [
        ("FAST-LIVO2 实时三维建图", "融合 LiDAR / Visual / IMU 信息，输出动态地图；当前阶段使用仿真点云验证规划算法，后续接入实时建图。", BLUE),
        ("MPC 局部轨迹生成", "面向异构机器人生成满足动力学约束、控制约束与安全距离要求的平滑轨迹。", GREEN),
        ("K-CBS 思想 + CB-MPC 框架", "兼顾全局冲突搜索与局部滚动优化，处理连续空间中的碰撞、死锁与执行偏差。", AMBER),
        ("RACER 分布式探索机制", "引入去中心化探索思路，提高扩展性，降低单点故障对全局任务的影响。", CYAN),
    ]
    for i, (title, body, color) in enumerate(rows):
        add_card(s, Cm(1.35), Cm(4.35 + i * 2.85), Cm(24.2), Cm(2.2), title, body, accent=color)
        if i < 3:
            add_arrow(s, Cm(26.0), Cm(5.45 + i * 2.85), Cm(26.0), Cm(6.05 + i * 2.85), color)
    add_metric(s, Cm(27.0), Cm(5.0), Cm(4.7), Cm(2.25), "安全", "显式约束", RED)
    add_metric(s, Cm(27.0), Cm(7.75), Cm(4.7), Cm(2.25), "实时", "滚动规划", BLUE)
    add_metric(s, Cm(27.0), Cm(10.5), Cm(4.7), Cm(2.25), "可扩展", "分布式协同", GREEN)
    add_footer(s)


def slide_innovation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "两项核心创新", "创新点")
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.45), Cm(3.75), Cm(14.85), Cm(10.4))
    left.adjustments[0] = 0.07
    set_fill(left, RGBColor(246, 250, 253))
    set_line(left)
    right = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(17.1), Cm(3.75), Cm(14.85), Cm(10.4))
    right.adjustments[0] = 0.07
    set_fill(right, RGBColor(248, 251, 248))
    set_line(right)
    add_text(s, Cm(2.05), Cm(4.35), Cm(13.6), Cm(0.8), "1. 实时“感知-规划”紧耦合", size=20, color=BLUE, bold=True)
    add_bullets(s, Cm(2.15), Cm(5.65), Cm(13.35), Cm(5.5), [
        "FAST-LIVO2 动态地图输出直接馈入 MPC 规划器。",
        "规划指令与感知数据闭环更新，减少串行流水线延迟。",
        "基于动态障碍物短期预测主动调整局部路径。"
    ], size=15)
    add_text(s, Cm(2.15), Cm(12.05), Cm(13.5), Cm(0.8), "价值：从“被动避障”走向“前瞻性避障”。", size=14.5, color=NAVY, bold=True)

    add_text(s, Cm(17.75), Cm(4.35), Cm(13.5), Cm(0.8), "2. 异构协同冲突解决范式", size=20, color=GREEN, bold=True)
    add_bullets(s, Cm(17.85), Cm(5.65), Cm(13.35), Cm(5.5), [
        "CB-MPC 兼顾全局路径最优性与局部实时调整。",
        "把异构平台运动学差异纳入冲突检测与消解。",
        "后续接入 RACER 后，可扩展至去中心化集群探索。"
    ], size=15)
    add_text(s, Cm(17.85), Cm(12.05), Cm(13.5), Cm(0.8), "价值：支持 UAV / UGV 精细化轨迹协同。", size=14.5, color=NAVY, bold=True)
    add_footer(s)


def slide_progress(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "当前进展：算法实现与仿真验证已打通", "项目进展")
    add_metric(s, Cm(1.55), Cm(3.55), Cm(6.2), Cm(2.6), "4架", "无人机对向飞行避碰", BLUE)
    add_metric(s, Cm(8.35), Cm(3.55), Cm(6.2), Cm(2.6), "80个", "随机障碍物环境", GREEN)
    add_metric(s, Cm(15.15), Cm(3.55), Cm(6.2), Cm(2.6), "2Hz", "持续重规划频率", AMBER)
    add_metric(s, Cm(21.95), Cm(3.55), Cm(6.2), Cm(2.6), "30Hz", "控制频率", RED)
    add_card(s, Cm(1.55), Cm(7.25), Cm(14.5), Cm(5.55), "规划组", "完成 CB-MPC 多机器人协同运动规划算法核心实现；将底层 MPC 控制器与多机器人冲突检测机制耦合；在 EGO-Planner 仿真环境中验证实时避碰。", BLUE)
    add_card(s, Cm(17.0), Cm(7.25), Cm(14.5), Cm(5.55), "建图组", "完成 CB-MPC 与 EGO-Planner 无人机仿真平台对接；搭建 ROS 话题通信链路，开发轨迹服务器，并集成 RViz 可视化。", GREEN)
    add_text(s, Cm(1.55), Cm(14.1), Cm(28.5), Cm(0.9), "阶段性结论：系统已形成从仿真环境、轨迹生成、冲突检测到可视化验证的基础闭环。", size=16, color=NAVY, bold=True)
    add_footer(s)


def slide_applications(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "落地应用：高风险、高复杂度、高频巡检场景", "应用场景")
    scenarios = [
        ("灾后搜救与应急救援", "无人机快速穿越障碍获取灾区全景；UGV 在引导下进入危险区域，执行探测与物资投放。", RED),
        ("大型工业园区智慧巡检", "无人机进行快速巡检和热成像检测；地面机器人针对异常点位进行精细复检和数据采集。", BLUE),
        ("城市治理与安防监控", "形成“空中宏观预警 + 地面微观精检”的现代化防控体系，服务大型活动、边防、交通管理。", GREEN),
    ]
    for i, (title, body, color) in enumerate(scenarios):
        add_card(s, Cm(1.35 + i * 10.45), Cm(4.05), Cm(9.45), Cm(7.4), title, body, accent=color)
    add_text(s, Cm(2.0), Cm(12.95), Cm(29.0), Cm(1.1), "共同需求：未知环境快速建图、动态障碍安全绕行、多机器人任务合理分工、弱通信条件下稳定执行。", size=17.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "市场前景：低空经济与机器人产业交汇", "市场数据")
    add_text(s, Cm(1.25), Cm(3.1), Cm(29.5), Cm(0.8), "备用文档显示，部分原始市场预测口径需谨慎引用；本页仅保留可查证或有多源印证的数据。", size=14.5, color=MUTED)

    # Bar chart: inspection robots.
    add_text(s, Cm(1.35), Cm(4.45), Cm(12.5), Cm(0.65), "中国智能巡检机器人市场规模", size=17, color=NAVY, bold=True)
    chart_x = Cm(1.6)
    chart_y = Cm(5.45)
    max_w = Cm(10.2)
    vals = [(2023, 19.71, BLUE), (2025, 27.48, GREEN)]
    for i, (year, val, color) in enumerate(vals):
        y = chart_y + Cm(i * 1.25)
        add_text(s, chart_x, y + Cm(0.08), Cm(1.45), Cm(0.45), str(year), size=11.5, color=MUTED)
        bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, chart_x + Cm(1.55), y, max_w, Cm(0.48))
        set_fill(bg, RGBColor(232, 238, 244))
        bg.line.fill.background()
        bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, chart_x + Cm(1.55), y, max_w * (val / 30), Cm(0.48))
        set_fill(bar, color)
        bar.line.fill.background()
        add_text(s, chart_x + Cm(1.75) + max_w * (val / 30), y - Cm(0.02), Cm(2.2), Cm(0.55), f"{val:.2f}亿元", size=12.5, color=INK, bold=True)
    add_text(s, Cm(1.6), Cm(8.15), Cm(12.5), Cm(0.8), "2025 年采用智研咨询/搜狐转载口径；备用文档提示“32亿元”未获直接印证。", size=9.5, color=MUTED)

    # Bar chart: drone market.
    add_text(s, Cm(17.0), Cm(4.45), Cm(12.5), Cm(0.65), "全球工业级无人机市场", size=17, color=NAVY, bold=True)
    vals2 = [(2025, 87.9, CYAN), (2035, 300, AMBER)]
    for i, (year, val, color) in enumerate(vals2):
        y = chart_y + Cm(i * 1.25)
        add_text(s, Cm(17.25), y + Cm(0.08), Cm(1.45), Cm(0.45), str(year), size=11.5, color=MUTED)
        bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(18.8), y, max_w, Cm(0.48))
        set_fill(bg, RGBColor(232, 238, 244))
        bg.line.fill.background()
        bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(18.8), y, max_w * (val / 320), Cm(0.48))
        set_fill(bar, color)
        bar.line.fill.background()
        add_text(s, Cm(19.0) + max_w * (val / 320), y - Cm(0.02), Cm(3.1), Cm(0.55), f"{val:g}亿美元", size=12.5, color=INK, bold=True)
    add_text(s, Cm(17.25), Cm(8.15), Cm(12.5), Cm(0.8), "WiseGuy Reports 预测：2025 年 87.9 亿美元，2035 年 300 亿美元。", size=9.5, color=MUTED)

    add_card(s, Cm(1.35), Cm(10.25), Cm(14.7), Cm(3.5), "产业融合窗口", "空地一体集群智慧 AI 场景解决方案融合低空经济与机器人两大万亿级产业。", BLUE)
    add_card(s, Cm(17.0), Cm(10.25), Cm(14.7), Cm(3.5), "应用市场空间", "范丛明公开发言：相关场景解决方案未来 5 年将催生超百亿元规模的应用市场。", GREEN)
    add_footer(s, "数据来源：机器人大讲堂(2024-10-14)；搜狐网转载智研咨询报告；WiseGuy Reports《Industrial Grade Drone Market Insights》；东方财富网/新浪财经(2025-10-31)；深圳市人工智能产业协会官方平台。")


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "团队分工：建图组 × 规划组", "团队介绍")
    add_card(s, Cm(1.35), Cm(3.7), Cm(14.7), Cm(6.6), "建图组", "张婧瑶：项目统筹、仿真环境搭建、ROS 系统集成、CB-MPC 与 EGO-Planner 对接\n张释文：建图模块数据处理与环境建模", BLUE)
    add_card(s, Cm(17.0), Cm(3.7), Cm(14.7), Cm(6.6), "规划组", "闫天路：CB-MPC 核心算法实现与调试\n陈可欣：规划算法设计与文献调研\n荣祎铭：规划算法搭建与仿真验证", GREEN)
    add_text(s, Cm(2.0), Cm(12.0), Cm(29.0), Cm(1.0), "团队能力覆盖 ROS、运动规划、三维感知、环境建模与仿真验证，已形成算法研究和系统集成协作闭环。", size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "后续规划与预期成果", "实施路径")
    steps = [
        ("2026.05\n已完成", "CB-MPC 核心实现\n4机仿真验证\nEGO-Planner 对接", BLUE),
        ("2026.06-08", "接入 FAST-LIVO2\n加入 UGV 模型\n动态障碍验证", GREEN),
        ("2026.09-11", "强化学习探索策略\n大规模仿真对比\n真实平台初测", AMBER),
        ("2026.11-2027.03", "室内/户外实物实验\n系统迭代优化\n成果完善", RED),
    ]
    x0 = Cm(1.45)
    for i, (time, body, color) in enumerate(steps):
        x = x0 + Cm(i * 7.8)
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Cm(2.15), Cm(4.0), Cm(1.1), Cm(1.1))
        set_fill(circ, color)
        circ.line.fill.background()
        if i < 3:
            line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Cm(3.25), Cm(4.52), Cm(5.65), Cm(0.08))
            set_fill(line, RGBColor(205, 217, 229))
            line.line.fill.background()
        add_text(s, x, Cm(5.55), Cm(5.6), Cm(1.25), time, size=15.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(7.25), Cm(5.6), Cm(2.6), body, size=12.2, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(s, Cm(1.5), Cm(12.0), Cm(14.7), Cm(3.3), "算法成果", "形成面向空地异构多机器人系统的统一环境建模、协同决策、冲突处理与自主探索方法。", BLUE)
    add_card(s, Cm(17.05), Cm(12.0), Cm(14.7), Cm(3.3), "系统成果", "完成仿真与典型真实场景实验验证，为巡检、环境探索等应用提供技术参考。", GREEN)
    add_footer(s)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    # remove default first slide if present by using blank layout only via new slides
    slide_cover(prs)
    slide_context(prs)
    slide_painpoints(prs)
    slide_architecture(prs)
    slide_algorithm(prs)
    slide_innovation(prs)
    slide_progress(prs)
    slide_applications(prs)
    slide_market(prs)
    slide_team(prs)
    slide_roadmap(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
